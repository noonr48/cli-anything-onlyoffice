#!/usr/bin/env python3
"""PPTX operations for the OnlyOffice CLI."""

from __future__ import annotations

import io
import os
import re
import tempfile
import warnings
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PPTXOperations:
    """Encapsulate PPTX creation, editing, inspection, and preview workflows."""

    MAX_EXTRACT_IMAGES = 1_000
    MAX_EXTRACT_IMAGE_COMPRESSED_BYTES = 50 * 1024 * 1024
    MAX_EXTRACT_IMAGE_PIXELS = 64_000_000
    MAX_EXTRACT_TOTAL_COMPRESSED_BYTES = 250 * 1024 * 1024
    MAX_EXTRACT_TOTAL_PIXELS = 512_000_000
    MAX_OUTPUT_PREFIX_LENGTH = 80
    IMAGE_OUTPUT_FORMATS = {
        "png": ("png", "PNG"),
        "jpg": ("jpg", "JPEG"),
        "jpeg": ("jpg", "JPEG"),
    }
    SAFE_PREFIX_RE = re.compile(r"^[A-Za-z0-9_. -]+$")

    _LAYOUT_MAP = {
        "title_only": 5,
        "content": 1,
        "blank": 6,
        "two_content": 3,
        "comparison": 4,
    }

    _ALIGN_MAP = {
        "left": PP_ALIGN.LEFT if PPTX_AVAILABLE else None,
        "center": PP_ALIGN.CENTER if PPTX_AVAILABLE else None,
        "right": PP_ALIGN.RIGHT if PPTX_AVAILABLE else None,
        "justify": PP_ALIGN.JUSTIFY if PPTX_AVAILABLE else None,
    }

    def __init__(self, host: Any):
        self.host = host

    @classmethod
    def normalize_image_format(cls, fmt: str) -> Any:
        fmt_key = str(fmt or "").strip().lower().lstrip(".")
        normalized = cls.IMAGE_OUTPUT_FORMATS.get(fmt_key)
        if normalized:
            return normalized, None
        return None, {
            "success": False,
            "error": "Unsupported image format. Use png or jpg.",
            "error_code": "unsupported_image_format",
            "allowed_formats": ["png", "jpg"],
        }

    @classmethod
    def validate_output_prefix(cls, prefix: str) -> Any:
        prefix_text = str(prefix or "").strip()
        preflight = {
            "operation": "pptx_extract_images",
            "prefix": prefix,
            "max_prefix_length": cls.MAX_OUTPUT_PREFIX_LENGTH,
        }
        if not prefix_text:
            preflight["status"] = "fail"
            preflight["reason"] = "prefix must not be empty"
        elif len(prefix_text) > cls.MAX_OUTPUT_PREFIX_LENGTH:
            preflight["status"] = "fail"
            preflight["reason"] = (
                f"prefix is {len(prefix_text)} characters; max is {cls.MAX_OUTPUT_PREFIX_LENGTH}"
            )
        elif (
            Path(prefix_text).is_absolute()
            or "/" in prefix_text
            or "\\" in prefix_text
            or ":" in prefix_text
            or prefix_text in {".", ".."}
            or any(part == ".." for part in Path(prefix_text).parts)
        ):
            preflight["status"] = "fail"
            preflight["reason"] = "prefix must be a filename prefix, not a path"
        elif not cls.SAFE_PREFIX_RE.match(prefix_text):
            preflight["status"] = "fail"
            preflight["reason"] = (
                "prefix may contain only letters, numbers, spaces, dots, hyphens, and underscores"
            )
        else:
            preflight["status"] = "pass"
            return prefix_text, preflight, None

        return None, preflight, {
            "success": False,
            "error": f"Unsafe PPTX image output prefix: {preflight['reason']}",
            "error_code": "unsafe_output_prefix",
            "preflight": preflight,
        }

    @staticmethod
    def prepare_output_dir(output_dir: str) -> Path:
        output_root = Path(output_dir).expanduser()
        output_root.mkdir(parents=True, exist_ok=True)
        return output_root.resolve()

    @staticmethod
    def safe_child_path(output_root: Path, filename: str) -> Path:
        candidate = (output_root / filename).resolve()
        try:
            common = os.path.commonpath([str(output_root), str(candidate)])
        except ValueError:
            common = ""
        if common != str(output_root):
            raise ValueError("Refusing to write outside output_dir")
        return candidate

    @classmethod
    def image_extract_resource_limits(cls) -> Dict[str, int]:
        return {
            "max_images": cls.MAX_EXTRACT_IMAGES,
            "max_compressed_image_bytes": cls.MAX_EXTRACT_IMAGE_COMPRESSED_BYTES,
            "max_decoded_image_pixels": cls.MAX_EXTRACT_IMAGE_PIXELS,
            "max_total_compressed_image_bytes": cls.MAX_EXTRACT_TOTAL_COMPRESSED_BYTES,
            "max_total_decoded_image_pixels": cls.MAX_EXTRACT_TOTAL_PIXELS,
        }

    @classmethod
    def load_bounded_image(
        cls,
        PILImage: Any,
        image_bytes: bytes,
        *,
        total_compressed_bytes: int = 0,
        total_decoded_pixels: int = 0,
    ) -> Any:
        compressed_size = len(image_bytes)
        compressed_after = total_compressed_bytes + compressed_size
        if compressed_after > cls.MAX_EXTRACT_TOTAL_COMPRESSED_BYTES:
            return None, None, {
                "skipped": True,
                "stop_extraction": True,
                "error": (
                    f"Embedded images would total {compressed_after} compressed bytes, exceeding "
                    f"the safe limit {cls.MAX_EXTRACT_TOTAL_COMPRESSED_BYTES}"
                ),
                "error_code": "aggregate_image_compressed_bytes_limit_exceeded",
                "compressed_size_bytes": compressed_size,
                "total_compressed_size_bytes": total_compressed_bytes,
                "would_total_compressed_size_bytes": compressed_after,
                "max_total_compressed_size_bytes": cls.MAX_EXTRACT_TOTAL_COMPRESSED_BYTES,
            }
        if compressed_size > cls.MAX_EXTRACT_IMAGE_COMPRESSED_BYTES:
            return None, None, {
                "skipped": True,
                "error": (
                    f"Embedded image is {compressed_size} compressed bytes, exceeding "
                    f"the safe limit {cls.MAX_EXTRACT_IMAGE_COMPRESSED_BYTES}"
                ),
                "error_code": "image_compressed_bytes_limit_exceeded",
                "compressed_size_bytes": compressed_size,
                "max_compressed_size_bytes": cls.MAX_EXTRACT_IMAGE_COMPRESSED_BYTES,
                "total_compressed_size_bytes": compressed_after,
            }

        img = None
        try:
            stream = io.BytesIO(image_bytes)
            bomb_warning = getattr(PILImage, "DecompressionBombWarning", None)
            with warnings.catch_warnings():
                if bomb_warning is not None:
                    warnings.simplefilter("error", bomb_warning)
                img = PILImage.open(stream)
                width, height = [int(value) for value in img.size]
                pixels = width * height
                pixels_after = total_decoded_pixels + pixels
                if pixels_after > cls.MAX_EXTRACT_TOTAL_PIXELS:
                    img.close()
                    return None, None, {
                        "skipped": True,
                        "stop_extraction": True,
                        "error": (
                            f"Embedded images would decode to {pixels_after} pixels, exceeding "
                            f"the safe limit {cls.MAX_EXTRACT_TOTAL_PIXELS}"
                        ),
                        "error_code": "aggregate_image_pixel_limit_exceeded",
                        "width": width,
                        "height": height,
                        "pixels": pixels,
                        "total_decoded_image_pixels": total_decoded_pixels,
                        "would_total_decoded_image_pixels": pixels_after,
                        "max_total_decoded_image_pixels": cls.MAX_EXTRACT_TOTAL_PIXELS,
                        "compressed_size_bytes": compressed_size,
                        "total_compressed_size_bytes": compressed_after,
                    }
                if pixels > cls.MAX_EXTRACT_IMAGE_PIXELS:
                    img.close()
                    return None, None, {
                        "skipped": True,
                        "error": (
                            f"Embedded image decodes to {pixels} pixels, exceeding "
                            f"the safe limit {cls.MAX_EXTRACT_IMAGE_PIXELS}"
                        ),
                        "error_code": "image_pixel_limit_exceeded",
                        "width": width,
                        "height": height,
                        "pixels": pixels,
                        "max_decoded_image_pixels": cls.MAX_EXTRACT_IMAGE_PIXELS,
                        "compressed_size_bytes": compressed_size,
                        "total_compressed_size_bytes": compressed_after,
                        "total_decoded_image_pixels": pixels_after,
                    }
                img.load()
            return img, {
                "width": width,
                "height": height,
                "pixels": pixels,
                "compressed_size_bytes": compressed_size,
                "total_compressed_size_bytes": compressed_after,
                "total_decoded_image_pixels": pixels_after,
            }, None
        except Exception as exc:
            if img is not None:
                img.close()
            return None, None, {
                "skipped": True,
                "error": f"Could not safely decode embedded image: {exc}",
                "error_code": "unsafe_image_decode",
                "compressed_size_bytes": compressed_size,
                "total_compressed_size_bytes": compressed_after,
            }

    @staticmethod
    def save_bounded_image(
        img: Any,
        out_path: Path,
        fmt_ext: str,
        pillow_format: str,
    ) -> None:
        tmp_path = out_path.with_name(f".{out_path.name}.tmp")
        try:
            if fmt_ext == "jpg":
                converted = img.convert("RGB")
                try:
                    converted.save(str(tmp_path), "JPEG", quality=90)
                finally:
                    converted.close()
            else:
                img.save(str(tmp_path), pillow_format)
            os.replace(str(tmp_path), str(out_path))
        except Exception:
            try:
                tmp_path.unlink()
            except FileNotFoundError:
                pass
            raise

    def create_presentation(
        self, output_path: str, title: str = "", subtitle: str = ""
    ) -> Dict[str, Any]:
        """Create a new presentation with a title slide."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            with self.host._file_lock(output_path):
                backup = self.host._snapshot_backup(output_path)
                prs = Presentation()
                prs.slide_width = Inches(13.333)  # 16:9 widescreen
                prs.slide_height = Inches(7.5)
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                if title:
                    slide.shapes.title.text = title
                if subtitle:
                    slide.placeholders[1].text = subtitle
                self.host._safe_save(prs, output_path)
            return {
                "success": True,
                "file": output_path,
                "title": title,
                "slides": 1,
                "size": Path(output_path).stat().st_size,
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def add_slide(
        self, file_path: str, title: str, content: str = "", layout: str = "content"
    ) -> Dict[str, Any]:
        """Add a slide using one of the supported preset layouts."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            with self.host._file_lock(file_path):
                backup = self.host._snapshot_backup(file_path)
                prs = Presentation(file_path)
                slide = prs.slides.add_slide(
                    prs.slide_layouts[self._LAYOUT_MAP.get(layout.lower(), 1)]
                )
                if slide.shapes.title:
                    slide.shapes.title.text = title
                if content and len(slide.placeholders) > 1:
                    slide.placeholders[1].text = content
                self.host._safe_save(prs, file_path)
            return {
                "success": True,
                "file": file_path,
                "title": title,
                "total_slides": len(prs.slides),
                "layout": layout,
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def add_bullet_slide(
        self, file_path: str, title: str, bullets: str
    ) -> Dict[str, Any]:
        """Add a title-and-content slide with bullet points separated by newlines."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            with self.host._file_lock(file_path):
                backup = self.host._snapshot_backup(file_path)
                prs = Presentation(file_path)
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                if slide.shapes.title:
                    slide.shapes.title.text = title
                bullet_lines = [bullet.strip() for bullet in bullets.split("\n") if bullet.strip()]
                if bullet_lines:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()
                    for bullet_index, bullet in enumerate(bullet_lines):
                        paragraph = text_frame.paragraphs[0] if bullet_index == 0 else text_frame.add_paragraph()
                        paragraph.text = bullet
                        paragraph.level = 0
                self.host._safe_save(prs, file_path)
            return {
                "success": True,
                "file": file_path,
                "title": title,
                "total_slides": len(prs.slides),
                "bullets": len(bullet_lines),
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def read_presentation(self, file_path: str) -> Dict[str, Any]:
        """Read all slides and visible text content from a presentation."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            prs = Presentation(file_path)
            slides_data = []
            for index, slide in enumerate(prs.slides, start=1):
                slide_info = {"slide_number": index, "title": "", "content": []}
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    is_title = shape == slide.shapes.title
                    if not is_title:
                        try:
                            placeholder = shape.placeholder_format
                            if placeholder is not None and placeholder.idx == 0:
                                is_title = True
                        except Exception:
                            pass
                    if is_title:
                        slide_info["title"] = text
                    else:
                        slide_info["content"].append(text)
                slides_data.append(slide_info)
            return {
                "success": True,
                "file": file_path,
                "slide_count": len(slides_data),
                "slides": slides_data,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def add_image_slide(
        self, file_path: str, title: str, image_path: str
    ) -> Dict[str, Any]:
        """Add a title-only slide and place an image onto it."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            with self.host._file_lock(file_path):
                backup = self.host._snapshot_backup(file_path)
                prs = Presentation(file_path)
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                if slide.shapes.title:
                    slide.shapes.title.text = title
                slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
                self.host._safe_save(prs, file_path)
            return {
                "success": True,
                "file": file_path,
                "title": title,
                "total_slides": len(prs.slides),
                "image": image_path,
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def add_table_slide(
        self,
        file_path: str,
        title: str,
        headers_csv: str,
        data_csv: str,
        coerce_rows: bool = False,
    ) -> Dict[str, Any]:
        """Add a title-only slide and populate it with a table."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            with self.host._file_lock(file_path):
                backup = self.host._snapshot_backup(file_path)
                prs = Presentation(file_path)
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                if slide.shapes.title:
                    slide.shapes.title.text = title

                headers = [header.strip() for header in headers_csv.split(",")]
                rows = []
                for row_str in data_csv.split(";"):
                    if row_str.strip():
                        rows.append([cell.strip() for cell in row_str.split(",")])

                ok, err, rows = self.host._validate_tabular_rows(
                    headers, rows, coerce_rows=coerce_rows
                )
                if not ok:
                    return {"success": False, "error": err}

                table_shape = slide.shapes.add_table(
                    len(rows) + 1,
                    len(headers),
                    Inches(0.5),
                    Inches(1.5),
                    Inches(9),
                    Inches(5),
                )
                table = table_shape.table
                for col_index, header in enumerate(headers):
                    table.cell(0, col_index).text = header
                for row_index, row_data in enumerate(rows, start=1):
                    for col_index, value in enumerate(row_data):
                        table.cell(row_index, col_index).text = value

                self.host._safe_save(prs, file_path)
            return {
                "success": True,
                "file": file_path,
                "title": title,
                "total_slides": len(prs.slides),
                "rows": len(rows),
                "columns": len(headers),
                "coerce_rows": bool(coerce_rows),
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def delete_slide(self, file_path: str, slide_index: int) -> Dict[str, Any]:
        """Delete a slide by 0-based index."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            with self.host._file_lock(file_path):
                backup = self.host._snapshot_backup(file_path)
                prs = Presentation(file_path)
                total = len(prs.slides)
                if slide_index < 0 or slide_index >= total:
                    return {
                        "success": False,
                        "error": f"Slide index {slide_index} out of range (0..{total - 1})",
                    }
                if total <= 1:
                    return {"success": False, "error": "Cannot delete the only slide"}
                rel_id = prs.slides._sldIdLst[slide_index].get(qn("r:id"))
                prs.part.drop_rel(rel_id)
                del prs.slides._sldIdLst[slide_index]
                self.host._safe_save(prs, file_path)
            return {
                "success": True,
                "file": file_path,
                "deleted_slide": slide_index,
                "remaining_slides": total - 1,
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def speaker_notes(
        self, file_path: str, slide_index: int, notes_text: str = None
    ) -> Dict[str, Any]:
        """Read or set speaker notes for a slide."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            if notes_text is None:
                prs = Presentation(file_path)
                total = len(prs.slides)
                if slide_index < 0 or slide_index >= total:
                    return {"success": False, "error": f"Slide index {slide_index} out of range"}
                slide = prs.slides[slide_index]
                notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                return {
                    "success": True,
                    "file": file_path,
                    "slide_index": slide_index,
                    "has_notes": slide.has_notes_slide,
                    "notes": notes,
                    "mode": "read",
                }

            with self.host._file_lock(file_path):
                backup = self.host._snapshot_backup(file_path)
                prs = Presentation(file_path)
                total = len(prs.slides)
                if slide_index < 0 or slide_index >= total:
                    return {"success": False, "error": f"Slide index {slide_index} out of range"}
                slide = prs.slides[slide_index]
                slide.notes_slide.notes_text_frame.text = notes_text
                self.host._safe_save(prs, file_path)
            return {
                "success": True,
                "file": file_path,
                "slide_index": slide_index,
                "notes": notes_text[:200],
                "mode": "write",
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def update_slide_text(
        self,
        file_path: str,
        slide_index: int,
        title: str = None,
        body: str = None,
    ) -> Dict[str, Any]:
        """Update the title and/or first body text frame on a slide."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            with self.host._file_lock(file_path):
                backup = self.host._snapshot_backup(file_path)
                prs = Presentation(file_path)
                total = len(prs.slides)
                if slide_index < 0 or slide_index >= total:
                    return {"success": False, "error": f"Slide index {slide_index} out of range"}
                slide = prs.slides[slide_index]
                updated_fields = []
                if title is not None and slide.shapes.title:
                    slide.shapes.title.text = title
                    updated_fields.append("title")
                if body is not None:
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape != slide.shapes.title:
                            shape.text_frame.text = body
                            updated_fields.append("body")
                            break
                self.host._safe_save(prs, file_path)
            return {
                "success": True,
                "file": file_path,
                "slide_index": slide_index,
                "updated_fields": updated_fields,
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def slide_count(self, file_path: str) -> Dict[str, Any]:
        """Return slide count and slide titles."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            prs = Presentation(file_path)
            titles = []
            for slide in prs.slides:
                titles.append(slide.shapes.title.text if slide.shapes.title else "")
            return {
                "success": True,
                "file": file_path,
                "slide_count": len(prs.slides),
                "slide_titles": titles,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def extract_images(
        self,
        file_path: str,
        output_dir: str,
        slide_index: int = None,
        fmt: str = "png",
        prefix: str = "slide",
    ) -> Dict[str, Any]:
        """Extract all images from a presentation, optionally limited to one slide."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            from PIL import Image as PILImage

            fmt_info, fmt_error = self.normalize_image_format(fmt)
            if fmt_error:
                return fmt_error
            fmt_ext, pillow_format = fmt_info
            safe_prefix, prefix_preflight, prefix_error = self.validate_output_prefix(prefix)
            if prefix_error:
                return prefix_error
            output_root = self.prepare_output_dir(output_dir)
            prs = Presentation(file_path)
            extracted = []
            index = 0
            truncated = False
            warnings = []
            resource_limits = self.image_extract_resource_limits()
            total_compressed_bytes = 0
            total_decoded_pixels = 0
            if slide_index is not None:
                if slide_index < 0 or slide_index >= len(prs.slides):
                    return {
                        "success": False,
                        "error": f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})",
                    }
                slides_to_scan = [(slide_index, prs.slides[slide_index])]
            else:
                slides_to_scan = list(enumerate(prs.slides))

            for slide_number, slide in slides_to_scan:
                for shape in slide.shapes:
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue
                    if index >= self.MAX_EXTRACT_IMAGES:
                        truncated = True
                        warning = f"Stopped after {self.MAX_EXTRACT_IMAGES} images"
                        if warning not in warnings:
                            warnings.append(warning)
                        break
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        out_path = self.safe_child_path(
                            output_root,
                            f"{safe_prefix}_{slide_number:02d}_{index:03d}.{fmt_ext}",
                        )
                        img, image_meta, skip_entry = self.load_bounded_image(
                            PILImage,
                            image_bytes,
                            total_compressed_bytes=total_compressed_bytes,
                            total_decoded_pixels=total_decoded_pixels,
                        )
                        stop_after_current = False
                        if skip_entry:
                            if (
                                skip_entry.get("error_code")
                                != "aggregate_image_compressed_bytes_limit_exceeded"
                            ):
                                total_compressed_bytes = skip_entry.get(
                                    "total_compressed_size_bytes",
                                    total_compressed_bytes,
                                )
                            if (
                                skip_entry.get("error_code")
                                != "aggregate_image_pixel_limit_exceeded"
                            ):
                                total_decoded_pixels = skip_entry.get(
                                    "total_decoded_image_pixels",
                                    total_decoded_pixels,
                                )
                            stop_after_current = bool(skip_entry.get("stop_extraction"))
                            if stop_after_current:
                                truncated = True
                                warning = str(skip_entry.get("error") or "Stopped at image extraction budget")
                                if warning not in warnings:
                                    warnings.append(warning)
                            skip_entry.update(
                                {
                                    "index": index,
                                    "slide": slide_number,
                                    "shape_name": shape.name,
                                }
                            )
                            extracted.append(skip_entry)
                        else:
                            total_compressed_bytes = image_meta["total_compressed_size_bytes"]
                            total_decoded_pixels = image_meta["total_decoded_image_pixels"]
                            try:
                                self.save_bounded_image(img, out_path, fmt_ext, pillow_format)
                            finally:
                                img.close()
                            extracted.append(
                                {
                                    "index": index,
                                    "slide": slide_number,
                                    "file": str(out_path),
                                    "format": fmt_ext,
                                    "width": image_meta["width"],
                                    "height": image_meta["height"],
                                    "pixels": image_meta["pixels"],
                                    "compressed_size_bytes": image_meta["compressed_size_bytes"],
                                    "size_bytes": os.path.getsize(str(out_path)),
                                    "shape_name": shape.name,
                                }
                            )
                        if stop_after_current:
                            index += 1
                            break
                    except Exception as img_err:
                        extracted.append(
                            {
                                "index": index,
                                "slide": slide_number,
                                "skipped": True,
                                "error": str(img_err),
                                "error_code": "image_extract_failed",
                            }
                        )
                    index += 1
                if truncated:
                    break
            return {
                "success": True,
                "file": file_path,
                "output_dir": str(output_root),
                "prefix": safe_prefix,
                "images_extracted": len([item for item in extracted if "file" in item]),
                "images_skipped": len([item for item in extracted if item.get("skipped")]),
                "images": extracted,
                "truncated": truncated,
                "warnings": warnings,
                "resource_limits": resource_limits,
                "resource_usage": {
                    "compressed_image_bytes": total_compressed_bytes,
                    "decoded_image_pixels": total_decoded_pixels,
                },
                "preflight": prefix_preflight,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def list_shapes(self, file_path: str, slide_index: int = None) -> Dict[str, Any]:
        """List shapes with spatial info, text, and placeholder metadata."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            prs = Presentation(file_path)
            slide_width_emu = prs.slide_width
            slide_height_emu = prs.slide_height
            if slide_index is not None:
                if slide_index < 0 or slide_index >= len(prs.slides):
                    return {
                        "success": False,
                        "error": f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})",
                    }
                slides_to_scan = [(slide_index, prs.slides[slide_index])]
            else:
                slides_to_scan = list(enumerate(prs.slides))

            slides = []
            for slide_number, slide in slides_to_scan:
                shapes = []
                for shape in slide.shapes:
                    info = {
                        "shape_id": shape.shape_id,
                        "name": shape.name,
                        "shape_type": str(shape.shape_type),
                        "left_inches": round(shape.left / 914400, 3) if shape.left is not None else None,
                        "top_inches": round(shape.top / 914400, 3) if shape.top is not None else None,
                        "width_inches": round(shape.width / 914400, 3) if shape.width is not None else None,
                        "height_inches": round(shape.height / 914400, 3) if shape.height is not None else None,
                        "rotation": shape.rotation,
                        "has_text": shape.has_text_frame,
                    }
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        info["text"] = text_frame.text[:200]
                        info["word_wrap"] = text_frame.word_wrap
                        info["auto_size"] = str(text_frame.auto_size) if text_frame.auto_size else None
                        info["margin_left"] = round(text_frame.margin_left / 914400, 3) if text_frame.margin_left is not None else None
                        info["margin_right"] = round(text_frame.margin_right / 914400, 3) if text_frame.margin_right is not None else None
                        info["margin_top"] = round(text_frame.margin_top / 914400, 3) if text_frame.margin_top is not None else None
                        info["margin_bottom"] = round(text_frame.margin_bottom / 914400, 3) if text_frame.margin_bottom is not None else None
                    try:
                        placeholder = shape.placeholder_format
                        if placeholder is not None:
                            info["placeholder_idx"] = placeholder.idx
                            info["placeholder_type"] = str(placeholder.type)
                    except Exception:
                        pass
                    if info["left_inches"] is not None and info["width_inches"] is not None:
                        info["right_inches"] = round(info["left_inches"] + info["width_inches"], 3)
                    if info["top_inches"] is not None and info["height_inches"] is not None:
                        info["bottom_inches"] = round(info["top_inches"] + info["height_inches"], 3)
                    shapes.append(info)
                slides.append(
                    {
                        "slide_index": slide_number,
                        "shape_count": len(shapes),
                        "shapes": shapes,
                    }
                )

            return {
                "success": True,
                "file": file_path,
                "slide_width_inches": round(slide_width_emu / 914400, 3),
                "slide_height_inches": round(slide_height_emu / 914400, 3),
                "slides": slides,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def add_textbox(
        self,
        file_path: str,
        slide_index: int,
        text: str,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 5.0,
        height: float = 1.5,
        font_size: float = None,
        font_name: str = None,
        bold: bool = False,
        italic: bool = False,
        color: str = None,
        align: str = None,
        word_wrap: bool = True,
    ) -> Dict[str, Any]:
        """Add a textbox at a specific position on a slide."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            with self.host._file_lock(file_path):
                backup = self.host._snapshot_backup(file_path)
                prs = Presentation(file_path)
                total = len(prs.slides)
                if slide_index < 0 or slide_index >= total:
                    return {
                        "success": False,
                        "error": f"Slide index {slide_index} out of range (0-{total-1})",
                    }
                slide = prs.slides[slide_index]
                textbox = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height),
                )
                text_frame = textbox.text_frame
                text_frame.word_wrap = word_wrap
                paragraph = text_frame.paragraphs[0]
                paragraph.text = text
                if align and align.lower() in self._ALIGN_MAP:
                    paragraph.alignment = self._ALIGN_MAP[align.lower()]
                for run in paragraph.runs:
                    if font_size is not None:
                        run.font.size = Pt(font_size)
                    if font_name:
                        run.font.name = font_name
                    if bold:
                        run.font.bold = True
                    if italic:
                        run.font.italic = True
                    if color:
                        hex_color = color.lstrip("#")
                        run.font.color.rgb = PptxRGBColor(
                            int(hex_color[:2], 16),
                            int(hex_color[2:4], 16),
                            int(hex_color[4:6], 16),
                        )
                self.host._safe_save(prs, file_path)
            return {
                "success": True,
                "file": file_path,
                "slide_index": slide_index,
                "shape_name": textbox.name,
                "position": {"left": left, "top": top, "width": width, "height": height},
                "text": text[:100],
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def modify_shape(
        self,
        file_path: str,
        slide_index: int,
        shape_name: str,
        left: float = None,
        top: float = None,
        width: float = None,
        height: float = None,
        text: str = None,
        font_size: float = None,
        rotation: float = None,
    ) -> Dict[str, Any]:
        """Modify a shape by name on a slide."""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx not installed"}
        try:
            with self.host._file_lock(file_path):
                backup = self.host._snapshot_backup(file_path)
                prs = Presentation(file_path)
                total = len(prs.slides)
                if slide_index < 0 or slide_index >= total:
                    return {
                        "success": False,
                        "error": f"Slide index {slide_index} out of range (0-{total-1})",
                    }
                slide = prs.slides[slide_index]
                target = None
                for shape in slide.shapes:
                    if shape.name == shape_name:
                        target = shape
                        break
                if target is None:
                    return {
                        "success": False,
                        "error": f"Shape '{shape_name}' not found. Available: {[shape.name for shape in slide.shapes]}",
                    }

                changes = []
                if left is not None:
                    target.left = Inches(left)
                    changes.append(f"left={left}in")
                if top is not None:
                    target.top = Inches(top)
                    changes.append(f"top={top}in")
                if width is not None:
                    target.width = Inches(width)
                    changes.append(f"width={width}in")
                if height is not None:
                    target.height = Inches(height)
                    changes.append(f"height={height}in")
                if rotation is not None:
                    target.rotation = rotation
                    changes.append(f"rotation={rotation}°")
                if text is not None and target.has_text_frame:
                    target.text_frame.text = text
                    changes.append(f"text='{text[:50]}'")
                if font_size is not None and target.has_text_frame:
                    for paragraph in target.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size)
                    changes.append(f"font_size={font_size}pt")
                self.host._safe_save(prs, file_path)
            return {
                "success": True,
                "file": file_path,
                "slide_index": slide_index,
                "shape_name": shape_name,
                "changes": changes,
                "backup": backup or None,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}

    def preview_slide(
        self,
        file_path: str,
        output_dir: str,
        slide_index: int = None,
        dpi: int = 150,
    ) -> Dict[str, Any]:
        """Render presentation slides as PNG images via the shared Office-to-PDF path."""
        try:
            import fitz
        except ImportError:
            return {"success": False, "error": "PyMuPDF (fitz) not installed. Run: pip install PyMuPDF"}

        pdf_path = None
        try:
            os.makedirs(output_dir, exist_ok=True)
            fd, pdf_path = tempfile.mkstemp(
                prefix=f".{Path(file_path).stem}_preview_",
                suffix=".pdf",
                dir=output_dir,
            )
            os.close(fd)

            pdf_result = self.host._office_to_pdf(file_path, output_path=pdf_path)
            if not pdf_result.get("success"):
                return pdf_result

            total_slides = pdf_result.get("pages")
            if total_slides is None:
                pdf = fitz.open(pdf_path)
                total_slides = len(pdf)
                pdf.close()

            if slide_index is not None and (slide_index < 0 or slide_index >= total_slides):
                return {
                    "success": False,
                    "error": f"Slide {slide_index} out of range (0-{total_slides-1})",
                }

            render_result = self.host.pdf_page_to_image(
                pdf_path,
                output_dir,
                pages=str(slide_index) if slide_index is not None else None,
                dpi=dpi,
                fmt="png",
            )
            if not render_result.get("success"):
                return render_result

            rendered = []
            for image in render_result.get("images", []):
                page_number = image.get("page")
                if page_number is None:
                    continue
                suffix = Path(image["file"]).suffix or ".png"
                slide_path = os.path.join(output_dir, f"slide_{page_number:03d}{suffix}")
                if image["file"] != slide_path:
                    os.replace(image["file"], slide_path)
                rendered.append(
                    {
                        "slide": page_number,
                        "file": slide_path,
                        "width": image.get("width"),
                        "height": image.get("height"),
                        "dpi": image.get("dpi", dpi),
                        "size_bytes": os.path.getsize(slide_path),
                    }
                )

            return {
                "success": True,
                "file": file_path,
                "output_dir": output_dir,
                "total_slides": total_slides,
                "slides_rendered": len(rendered),
                "dpi": dpi,
                "images": rendered,
            }
        except Exception as exc:
            return {"success": False, "error": str(exc)}
        finally:
            if pdf_path:
                try:
                    os.unlink(pdf_path)
                except OSError:
                    pass
