#!/usr/bin/env python3
"""
OnlyOffice Document Server API Client - Advanced Formatting Edition

Provides full programmatic control with advanced formatting:
- Documents (.docx) - Create, read, edit, format, highlight, comments
- Spreadsheets (.xlsx) - Create, read, edit, formulas, styling
"""

import requests
import json
import os
import shutil
import tempfile
import fcntl
import re
import hashlib
import threading
import xml.etree.ElementTree as ET
from pathlib import Path
from datetime import datetime, timezone
from typing import Dict, Any, List, Optional, Tuple
from contextlib import contextmanager
from zipfile import ZipFile, ZIP_DEFLATED

from cli_anything.onlyoffice.utils.doc_ops import DocumentOperations
from cli_anything.onlyoffice.utils.pdf_ops import PDFOperations
from cli_anything.onlyoffice.utils.pptx_ops import PPTXOperations
from cli_anything.onlyoffice.utils.rdf_ops import RDFOperations
from cli_anything.onlyoffice.utils.xlsx_ops import XLSXOperations

# Module-level per-path thread locks — serialises same-process threads before
# flock takes over for cross-process serialisation.
_thread_lock_map: dict[str, threading.Lock] = {}
_thread_lock_map_mu: threading.Lock = threading.Lock()

# Import document libraries with advanced formatting
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Mm
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
    from docx.enum.section import WD_ORIENT
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import Workbook, load_workbook
    from openpyxl.chart import (
        BarChart,
        LineChart,
        PieChart,
        ScatterChart,
        Reference,
    )
    from openpyxl.chart.label import DataLabelList
    from openpyxl.drawing.image import Image as XLImage
    import openpyxl.utils

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from scipy import stats as scipy_stats

    SCIPY_AVAILABLE = True
except ImportError:
    SCIPY_AVAILABLE = False


OOXML_NS = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
    "dc": "http://purl.org/dc/elements/1.1/",
    "dcterms": "http://purl.org/dc/terms/",
    "xsi": "http://www.w3.org/2001/XMLSchema-instance",
    "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

DOCX_PAGE_SIZES = {
    "a4": ("A4", (210.0, 297.0)),
    "letter": ("Letter", (215.9, 279.4)),
}


class DocumentServerClient:
    """Enhanced client for full document and spreadsheet control"""

    def __init__(
        self, server_url="http://localhost:8080", secret="sloane-os-secret-key-2026"
    ):
        self.server_url = server_url.rstrip("/")
        self.secret = secret
        self.headers = {"Authorization": f"Bearer {secret}"}
        self.backup_dir = Path(
            os.environ.get("ONLYOFFICE_BACKUP_DIR", "/tmp/sloane-onlyoffice-backups")
        )
        self.backup_dir.mkdir(parents=True, exist_ok=True)
        self._lock_timeout_seconds = int(
            os.environ.get("ONLYOFFICE_LOCK_TIMEOUT", "15")
        )
        self._doc_ops = DocumentOperations(self)
        self._pdf_ops = PDFOperations(self)
        self._pptx_ops = PPTXOperations(self)
        self._rdf_ops = RDFOperations(self)
        self._xlsx_ops = XLSXOperations(self)
        self.supported_formula_functions = {
            "SUM",
            "AVERAGE",
            "AVG",
            "MIN",
            "MAX",
            "ABS",
            "ROUND",
            "COUNT",
            "COUNTA",
        }

    def _file_key(self, file_path: str) -> str:
        abs_path = str(Path(file_path).resolve())
        return hashlib.sha1(abs_path.encode("utf-8")).hexdigest()[:10]

    def _lock_path(self, file_path: str) -> Path:
        abs_path = str(Path(file_path).resolve())
        return self.backup_dir / f"{self._file_key(abs_path)}.lock"

    @contextmanager
    def _file_lock(self, file_path: str):
        """Two-layer file lock: threading.Lock (intra-process) + fcntl.flock (cross-process).

        fcntl.flock is per-process on Linux — threads within the same process bypass it.
        The threading.Lock layer serialises concurrent threads before flock serialises
        concurrent processes.
        """
        import time
        abs_path = str(Path(file_path).resolve())
        # Retrieve or create the per-path thread lock
        with _thread_lock_map_mu:
            if abs_path not in _thread_lock_map:
                _thread_lock_map[abs_path] = threading.Lock()
            tlock = _thread_lock_map[abs_path]

        with tlock:  # serialise threads within this process
            lock_path = str(self._lock_path(abs_path))
            lock_file = open(lock_path, "w")
            deadline = time.monotonic() + self._lock_timeout_seconds
            acquired = False
            try:
                while True:
                    try:
                        fcntl.flock(lock_file.fileno(), fcntl.LOCK_EX | fcntl.LOCK_NB)
                        acquired = True
                        break
                    except BlockingIOError:
                        if time.monotonic() >= deadline:
                            raise TimeoutError(
                                f"Could not acquire lock on {file_path} within "
                                f"{self._lock_timeout_seconds}s"
                            )
                        time.sleep(0.05)
                yield
            finally:
                if acquired:
                    try:
                        fcntl.flock(lock_file.fileno(), fcntl.LOCK_UN)
                    except Exception:
                        pass
                lock_file.close()
                try:
                    os.unlink(lock_path)
                except OSError:
                    pass

    @staticmethod
    def _current_orientation(section) -> str:
        return "landscape" if section.page_width > section.page_height else "portrait"

    @staticmethod
    def _custom_xml_summary(names: List[str]) -> Dict[str, List[str]]:
        all_parts = sorted(name for name in names if name.startswith("customXml/"))
        payload_parts = sorted(
            name
            for name in all_parts
            if name.count("/") == 1
            and name.endswith(".xml")
            and not Path(name).name.startswith("itemProps")
        )
        support_parts = sorted(name for name in all_parts if name not in payload_parts)
        return {
            "all_parts": all_parts,
            "payload_parts": payload_parts,
            "support_parts": support_parts,
        }

    @staticmethod
    def _atomic_zip_write(files: Dict[str, bytes], target_path: str) -> None:
        target = Path(target_path)
        target.parent.mkdir(parents=True, exist_ok=True)
        fd, tmp_path = tempfile.mkstemp(
            prefix=f".{target.name}.", suffix=".tmp", dir=str(target.parent)
        )
        os.close(fd)
        try:
            with ZipFile(tmp_path, "w", compression=ZIP_DEFLATED) as zout:
                for name, data in files.items():
                    zout.writestr(name, data)
            os.replace(tmp_path, str(target))
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    @staticmethod
    def _resolve_page_size(page_size: Optional[str]) -> Optional[Tuple[str, int, int]]:
        if not page_size:
            return None
        key = str(page_size).strip().lower()
        if key not in DOCX_PAGE_SIZES:
            raise ValueError(
                f"Unsupported page size '{page_size}'. Supported: {', '.join(sorted(name.upper() for name in DOCX_PAGE_SIZES))}"
            )
        label, (width_mm, height_mm) = DOCX_PAGE_SIZES[key]
        return label, Mm(width_mm), Mm(height_mm)

    @staticmethod
    def _label_page_size(width_emu: int, height_emu: int) -> str:
        if not width_emu or not height_emu:
            return "Unknown"
        width_mm = round(float(width_emu) / 36000.0, 1)
        height_mm = round(float(height_emu) / 36000.0, 1)
        dims = sorted([width_mm, height_mm])
        for label, (expected_w, expected_h) in DOCX_PAGE_SIZES.values():
            expected_dims = sorted([expected_w, expected_h])
            if all(abs(a - b) <= 1.5 for a, b in zip(dims, expected_dims)):
                return label
        return f"{width_mm}x{height_mm} mm"

    @staticmethod
    def _label_pdf_page_size(width_pts: float, height_pts: float) -> str:
        width_mm = round(float(width_pts) * 25.4 / 72.0, 1)
        height_mm = round(float(height_pts) * 25.4 / 72.0, 1)
        dims = sorted([width_mm, height_mm])
        for label, (expected_w, expected_h) in DOCX_PAGE_SIZES.values():
            expected_dims = sorted([expected_w, expected_h])
            if all(abs(a - b) <= 1.5 for a, b in zip(dims, expected_dims)):
                return label
        return f"{width_mm}x{height_mm} mm"

    @staticmethod
    def _normalized_font_name(name: Optional[str]) -> str:
        text = str(name or "").strip()
        return re.sub(r"\s+", " ", text).lower()

    def _document_sections_summary(self, doc) -> List[Dict[str, Any]]:
        sections = []
        for idx, section in enumerate(doc.sections):
            page_width_in = round(section.page_width.inches, 3)
            page_height_in = round(section.page_height.inches, 3)
            usable_width_in = round(
                section.page_width.inches
                - section.left_margin.inches
                - section.right_margin.inches,
                3,
            )
            usable_height_in = round(
                section.page_height.inches
                - section.top_margin.inches
                - section.bottom_margin.inches,
                3,
            )
            sections.append(
                {
                    "section_index": idx,
                    "page_size": self._label_page_size(
                        section.page_width, section.page_height
                    ),
                    "orientation": self._current_orientation(section),
                    "page_width_in": page_width_in,
                    "page_height_in": page_height_in,
                    "usable_width_in": usable_width_in,
                    "usable_height_in": usable_height_in,
                    "margins_in": {
                        "top": round(section.top_margin.inches, 3),
                        "bottom": round(section.bottom_margin.inches, 3),
                        "left": round(section.left_margin.inches, 3),
                        "right": round(section.right_margin.inches, 3),
                    },
                }
            )
        return sections

    def _collect_text_runs(self, doc) -> List[Dict[str, Any]]:
        runs_payload = []
        normal_style = doc.styles["Normal"] if "Normal" in [s.name for s in doc.styles] else None

        def iter_paragraphs():
            for idx, para in enumerate(doc.paragraphs):
                yield ("paragraph", idx, para)
            for ti, table in enumerate(doc.tables):
                for ri, row in enumerate(table.rows):
                    for ci, cell in enumerate(row.cells):
                        for pi, para in enumerate(cell.paragraphs):
                            yield (f"table_{ti}_r{ri}_c{ci}", pi, para)

        for location, para_idx, para in iter_paragraphs():
            style_font = para.style.font if para.style is not None else None
            for run_idx, run in enumerate(para.runs):
                text = str(run.text or "")
                if not text.strip():
                    continue
                font_name = (
                    run.font.name
                    or (style_font.name if style_font is not None else None)
                    or (normal_style.font.name if normal_style is not None else None)
                )
                font_size = (
                    (run.font.size.pt if run.font.size else None)
                    or (
                        style_font.size.pt
                        if style_font is not None and style_font.size is not None
                        else None
                    )
                    or (
                        normal_style.font.size.pt
                        if normal_style is not None and normal_style.font.size is not None
                        else None
                    )
                )
                runs_payload.append(
                    {
                        "location": location,
                        "paragraph_index": para_idx,
                        "run_index": run_idx,
                        "text_preview": text[:60],
                        "font_name": font_name,
                        "font_name_normalized": self._normalized_font_name(font_name),
                        "font_size": round(float(font_size), 2) if font_size else None,
                        "bold": bool(run.bold) if run.bold is not None else False,
                        "italic": bool(run.italic) if run.italic is not None else False,
                    }
                )
        return runs_payload

    @staticmethod
    def _docx_story_xml_parts(names: List[str]) -> List[str]:
        return sorted(
            [
                name
                for name in names
                if name == "word/document.xml"
                or re.fullmatch(r"word/header\d+\.xml", name)
                or re.fullmatch(r"word/footer\d+\.xml", name)
                or name in {"word/footnotes.xml", "word/endnotes.xml"}
            ]
        )

    @staticmethod
    def _xml_local_name(tag: str) -> str:
        return tag.rsplit("}", 1)[-1] if "}" in tag else tag

    def _rewrite_story_tree(
        self,
        elem,
        *,
        remove_comment_nodes: bool = False,
        accept_revisions: bool = False,
        stats: Optional[Dict[str, int]] = None,
    ) -> None:
        if stats is None:
            stats = {}

        comment_nodes = {"commentRangeStart", "commentRangeEnd", "commentReference"}
        unwrap_nodes = {"ins", "moveTo"} if accept_revisions else set()
        drop_nodes = (
            {
                "del",
                "moveFrom",
                "moveFromRangeStart",
                "moveFromRangeEnd",
                "moveToRangeStart",
                "moveToRangeEnd",
                "customXmlDelRangeStart",
                "customXmlDelRangeEnd",
                "customXmlInsRangeStart",
                "customXmlInsRangeEnd",
                "customXmlMoveFromRangeStart",
                "customXmlMoveFromRangeEnd",
                "customXmlMoveToRangeStart",
                "customXmlMoveToRangeEnd",
                "cellDel",
                "cellMerge",
            }
            if accept_revisions
            else set()
        )

        new_children = []
        for child in list(elem):
            self._rewrite_story_tree(
                child,
                remove_comment_nodes=remove_comment_nodes,
                accept_revisions=accept_revisions,
                stats=stats,
            )
            local = self._xml_local_name(child.tag)
            if remove_comment_nodes and local in comment_nodes:
                stats["comment_nodes_removed"] = stats.get("comment_nodes_removed", 0) + 1
                continue
            if accept_revisions and local in unwrap_nodes:
                stats[f"{local}_accepted"] = stats.get(f"{local}_accepted", 0) + 1
                new_children.extend(list(child))
                continue
            if accept_revisions and local in drop_nodes:
                stats[f"{local}_removed"] = stats.get(f"{local}_removed", 0) + 1
                continue
            new_children.append(child)
        elem[:] = new_children

    @staticmethod
    def _strip_docx_relationship_targets(files: Dict[str, bytes], predicate) -> None:
        rel_tag = f"{{{OOXML_NS['rel']}}}Relationship"
        for name, blob in list(files.items()):
            if not name.endswith(".rels"):
                continue
            try:
                root = ET.fromstring(blob)
            except ET.ParseError:
                continue
            changed = False
            kept = []
            for child in list(root):
                target = child.get("Target", "")
                rel_type = child.get("Type", "")
                if child.tag == rel_tag and predicate(target, rel_type):
                    changed = True
                    continue
                kept.append(child)
            if changed:
                root[:] = kept
                files[name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    @staticmethod
    def _strip_docx_content_types(files: Dict[str, bytes], predicate) -> None:
        name = "[Content_Types].xml"
        if name not in files:
            return
        root = ET.fromstring(files[name])
        override_tag = f"{{{OOXML_NS['ct']}}}Override"
        default_tag = f"{{{OOXML_NS['ct']}}}Default"
        changed = False
        kept = []
        for child in list(root):
            if child.tag in {override_tag, default_tag} and predicate(child):
                changed = True
                continue
            kept.append(child)
        if changed:
            root[:] = kept
            files[name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    def _get_sheet(self, wb, sheet_name: str):
        """Get a worksheet with a friendly error if not found."""
        if sheet_name in wb.sheetnames:
            return wb[sheet_name]
        raise KeyError(
            f"Sheet '{sheet_name}' not found. Available sheets: {wb.sheetnames}. "
            f"Use --sheet {wb.sheetnames[0]} to target the first sheet."
        )

    def _snapshot_backup(self, file_path: str) -> str:
        """Create a point-in-time backup snapshot before mutating existing files."""
        src = Path(file_path)
        if not src.exists():
            return ""
        stamp = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%S%fZ")
        backup_name = f"{src.name}.{self._file_key(file_path)}.{stamp}.bak{src.suffix}"
        dst = self.backup_dir / backup_name
        shutil.copy2(src, dst)
        return str(dst)

    def _safe_save(self, wb_or_doc, file_path: str):
        """Atomic save to reduce risk of partial writes."""
        target = Path(file_path)
        fd, tmp_path = tempfile.mkstemp(
            prefix=f".{target.name}.", dir=str(target.parent)
        )
        os.close(fd)
        try:
            wb_or_doc.save(tmp_path)
            os.replace(tmp_path, str(target))
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    def _validate_tabular_rows(
        self,
        headers: List[Any],
        rows: List[List[Any]],
        coerce_rows: bool = False,
    ):
        expected_cols = len(headers)
        if expected_cols == 0:
            return False, "Headers cannot be empty", []

        normalized_rows = []
        for idx, row in enumerate(rows, 1):
            row_len = len(row)
            if row_len != expected_cols:
                if not coerce_rows:
                    return (
                        False,
                        f"Row {idx} has {row_len} columns, expected {expected_cols}. Use --coerce-rows to pad/truncate.",
                        [],
                    )
                if row_len < expected_cols:
                    row = row + [""] * (expected_cols - row_len)
                else:
                    row = row[:expected_cols]
            normalized_rows.append(row)

        return True, None, normalized_rows

    def _backup_candidates(self, file_path: str):
        p = Path(file_path)
        key = self._file_key(file_path)
        suffix = p.suffix
        candidates = []
        for bf in self.backup_dir.glob(f"{p.name}*.bak{suffix}"):
            name = bf.name
            if f".{key}." in name or re.match(
                rf"^{re.escape(p.name)}\.\d{{8}}T\d{{12}}Z\.bak{re.escape(suffix)}$",
                name,
            ):
                candidates.append(bf)
        return sorted(candidates, key=lambda x: x.stat().st_mtime, reverse=True)

    def list_backups(self, file_path: str, limit: int = 20) -> Dict[str, Any]:
        try:
            items = []
            for bf in self._backup_candidates(file_path)[: max(1, int(limit))]:
                st = bf.stat()
                items.append(
                    {
                        "backup": str(bf),
                        "name": bf.name,
                        "size": st.st_size,
                        "modified": datetime.fromtimestamp(
                            st.st_mtime, timezone.utc
                        ).isoformat(),
                    }
                )
            return {
                "success": True,
                "file": file_path,
                "count": len(items),
                "backups": items,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def prune_backups(
        self,
        file_path: str = None,
        keep: int = 20,
        older_than_days: int = None,
    ) -> Dict[str, Any]:
        try:
            keep = max(0, int(keep))
            now = datetime.now(timezone.utc).timestamp()
            if file_path:
                groups = {str(Path(file_path)): self._backup_candidates(file_path)}
            else:
                groups = {}
                for bf in self.backup_dir.glob("*.bak*"):
                    name = bf.name
                    parts = name.split(".")
                    if len(parts) < 4:
                        group = name
                    elif re.fullmatch(r"[a-f0-9]{10}", parts[-4]):
                        group = f"{'.'.join(parts[:-4])}.{parts[-4]}.{parts[-1]}"
                    else:
                        group = f"{'.'.join(parts[:-3])}.{parts[-1]}"
                    groups.setdefault(group, []).append(bf)

            deleted = []
            scanned = 0
            for files in groups.values():
                files = sorted(files, key=lambda x: x.stat().st_mtime, reverse=True)
                scanned += len(files)
                for idx, bf in enumerate(files):
                    st = bf.stat()
                    too_old = False
                    if older_than_days is not None:
                        too_old = (now - st.st_mtime) > (int(older_than_days) * 86400)
                    over_keep = idx >= keep
                    if too_old or over_keep:
                        try:
                            bf.unlink()
                            deleted.append(str(bf))
                        except Exception:
                            pass

            return {
                "success": True,
                "scanned": scanned,
                "deleted": len(deleted),
                "deleted_files": deleted,
                "keep": keep,
                "older_than_days": older_than_days,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def restore_backup(
        self,
        file_path: str,
        backup: str = None,
        latest: bool = False,
        dry_run: bool = False,
    ) -> Dict[str, Any]:
        try:
            target = Path(file_path)
            if backup:
                chosen = Path(backup)
                if not chosen.is_absolute():
                    chosen = self.backup_dir / backup
                if not chosen.exists():
                    return {"success": False, "error": f"Backup not found: {backup}"}
            else:
                backups = self._backup_candidates(file_path)
                if not backups:
                    return {"success": False, "error": "No backups found for file"}
                chosen = backups[0] if latest or not backup else None

            if dry_run:
                return {
                    "success": True,
                    "file": file_path,
                    "restore_from": str(chosen),
                    "dry_run": True,
                }

            with self._file_lock(file_path):
                pre_restore_backup = (
                    self._snapshot_backup(file_path) if target.exists() else ""
                )
                fd, tmp_path = tempfile.mkstemp(
                    prefix=f".{target.name}.restore.", dir=str(target.parent)
                )
                os.close(fd)
                try:
                    shutil.copy2(chosen, tmp_path)
                    os.replace(tmp_path, str(target))
                finally:
                    if os.path.exists(tmp_path):
                        os.unlink(tmp_path)

            return {
                "success": True,
                "file": file_path,
                "restored_from": str(chosen),
                "pre_restore_backup": pre_restore_backup or None,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ==================== BASIC DOCUMENT OPERATIONS ====================

    def create_document(
        self, output_path: str, title: str = "", content: str = ""
    ) -> Dict[str, Any]:
        return self._doc_ops.create_document(output_path, title, content)

    def read_document(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.read_document(file_path)

    def append_to_document(self, file_path: str, content: str) -> Dict[str, Any]:
        return self._doc_ops.append_to_document(file_path, content)

    def _replace_across_runs(self, paragraph, search_text: str, replace_text: str) -> int:
        """Replace text that may span multiple runs while preserving run-level formatting.
        Uses runs-only text universe (ignores tracked changes, hyperlinks, field codes).
        Collects all match positions first, then applies back-to-front to avoid infinite loops."""
        runs = paragraph.runs
        if not runs:
            return 0
        # Build text and char_map from runs only (not paragraph.text which includes non-run nodes)
        full = "".join(r.text for r in runs)
        if search_text not in full:
            return 0

        # Collect ALL match positions on the original string first (avoids infinite loop
        # when replace_text contains search_text)
        matches = []
        start = 0
        while True:
            idx = full.find(search_text, start)
            if idx == -1:
                break
            matches.append(idx)
            start = idx + len(search_text)  # non-overlapping

        if not matches:
            return 0

        # Apply replacements back-to-front so earlier positions stay valid
        for idx in reversed(matches):
            # Rebuild char_map fresh for current run state
            char_map = []
            for ri, run in enumerate(runs):
                for ci in range(len(run.text)):
                    char_map.append((ri, ci))

            end_idx = idx + len(search_text)
            if end_idx > len(char_map) or idx >= len(char_map):
                continue  # skip: match extends beyond run-text (tracked change / hyperlink region)

            first_run, first_offset = char_map[idx]
            last_run, last_offset = char_map[end_idx - 1]

            if first_run == last_run:
                r = runs[first_run]
                r.text = r.text[:first_offset] + replace_text + r.text[last_offset + 1:]
            else:
                fr = runs[first_run]
                fr.text = fr.text[:first_offset] + replace_text
                for mi in range(first_run + 1, last_run):
                    runs[mi].text = ""
                lr = runs[last_run]
                lr.text = lr.text[last_offset + 1:]

        return len(matches)

    def search_replace_document(
        self, file_path: str, search_text: str, replace_text: str
    ) -> Dict[str, Any]:
        return self._doc_ops.search_replace_document(file_path, search_text, replace_text)

    # ==================== ADVANCED FORMATTING ====================

    def format_paragraph(
        self,
        file_path: str,
        paragraph_index: int,
        bold: bool = False,
        italic: bool = False,
        underline: bool = False,
        font_name: str = None,
        font_size: int = None,
        color: str = None,
        alignment: str = None,
    ) -> Dict[str, Any]:
        return self._doc_ops.format_paragraph(
            file_path,
            paragraph_index,
            bold,
            italic,
            underline,
            font_name,
            font_size,
            color,
            alignment,
        )

    def highlight_text(
        self, file_path: str, search_text: str, color: str = "yellow"
    ) -> Dict[str, Any]:
        return self._doc_ops.highlight_text(file_path, search_text, color)

    def add_table(
        self, file_path: str, headers_csv: str, data_csv: str
    ) -> Dict[str, Any]:
        return self._doc_ops.add_table(file_path, headers_csv, data_csv)

    def add_comment(
        self, file_path: str, comment_text: str, paragraph_index: int = 0,
        author: str = "SLOANE Agent"
    ) -> Dict[str, Any]:
        return self._doc_ops.add_comment(
            file_path, comment_text, paragraph_index, author
        )

    # ==================== APA REFERENCE MANAGEMENT ====================

    def _apa_in_text(self, author_str: str, year: str) -> str:
        """Generate APA 7th in-text citation from author string.
        Handles: single author, two authors (& separator), 3+ authors (et al.)"""
        # Detect multiple authors by common separators
        # Forms: "Smith, J., Jones, A." or "Smith, J., & Jones, A." or "Smith and Jones"
        parts = author_str.strip()
        # Split by " & " or " and " first
        if " & " in parts or " and " in parts:
            authors = [a.strip() for a in parts.replace(" and ", " & ").split(" & ")]
            surnames = [a.split(",")[0].strip() for a in authors]
            if len(surnames) == 2:
                return f"({surnames[0]} & {surnames[1]}, {year})"
            elif len(surnames) >= 3:
                return f"({surnames[0]} et al., {year})"
        # Check for comma-separated authors: "Smith, J., Jones, A."
        # Pattern: if there are 3+ commas, likely multiple authors (surname, initial, surname, initial)
        commas = parts.count(",")
        if commas >= 3:
            # Multiple authors in "Surname, I., Surname, I." format
            # Split by ", " but keep pairs together
            tokens = [t.strip() for t in parts.split(",")]
            surnames = [tokens[i] for i in range(0, len(tokens), 2) if i < len(tokens)]
            if len(surnames) == 2:
                return f"({surnames[0]} & {surnames[1]}, {year})"
            elif len(surnames) >= 3:
                return f"({surnames[0]} et al., {year})"
        # Single author
        surname = parts.split(",")[0].strip()
        return f"({surname}, {year})"

    def add_reference(
        self, file_path: str, ref_json: str
    ) -> Dict[str, Any]:
        return self._doc_ops.add_reference(file_path, ref_json)

    def _format_apa7_reference(self, ref: dict):
        """Format a single reference in APA 7th edition style.
        Returns a list of {"text": str, "italic": bool} spans for proper Word rendering."""
        rtype = ref.get("type", "journal")
        author = ref.get("author", "Unknown")
        year = ref.get("year", "n.d.")
        title = ref.get("title", "Untitled")
        S = lambda t, i=False: {"text": t, "italic": i}  # span helper

        if rtype == "journal":
            source = ref.get("source", "")
            vol = ref.get("volume", "")
            issue = ref.get("issue", "")
            pages = ref.get("pages", "")
            doi = ref.get("doi", "")
            spans = [S(f"{author} ({year}). {title}. ")]
            if source:
                spans.append(S(source, True))  # journal name italic
                if vol:
                    spans.append(S(", "))
                    spans.append(S(vol, True))  # volume number italic
                    if issue:
                        spans.append(S(f"({issue})"))
                elif issue:
                    spans.append(S(f"({issue})"))
                if pages:
                    spans.append(S(f", {pages}"))
                spans.append(S(". "))
            if doi:
                url = f"https://doi.org/{doi}" if not doi.startswith("http") else doi
                spans.append(S(url))
            return spans

        elif rtype == "book":
            publisher = ref.get("source", ref.get("publisher", ""))
            doi = ref.get("doi", "")
            spans = [S(f"{author} ({year}). "), S(title, True), S(". ")]
            if publisher:
                spans.append(S(f"{publisher}. "))
            if doi:
                url = f"https://doi.org/{doi}" if not doi.startswith("http") else doi
                spans.append(S(url))
            return spans

        elif rtype == "chapter":
            editor = ref.get("editor", "")
            book_title = ref.get("source", "")
            pages = ref.get("pages", "")
            publisher = ref.get("publisher", "")
            spans = [S(f"{author} ({year}). {title}. ")]
            if editor:
                spans.append(S(f"In {editor} (Ed.), "))
            if book_title:
                spans.append(S(book_title, True))  # book title italic
            if pages:
                spans.append(S(f" (pp. {pages})"))
            spans.append(S(". "))
            if publisher:
                spans.append(S(f"{publisher}."))
            return spans

        elif rtype == "website":
            url = ref.get("url", ref.get("doi", ""))
            source = ref.get("source", "")
            spans = [S(f"{author} ({year}). "), S(title, True), S(". ")]
            if source:
                spans.append(S(f"{source}. "))
            if url:
                spans.append(S(url))
            return spans

        elif rtype == "report":
            source = ref.get("source", "")
            url = ref.get("url", "")
            spans = [S(f"{author} ({year}). "), S(title, True), S(". ")]
            if source:
                spans.append(S(f"{source}. "))
            if url:
                spans.append(S(url))
            return spans

        else:
            # Fallback: generic
            return [S(f"{author} ({year}). {title}.")]

    def build_references(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.build_references(file_path)

    def add_image(
        self,
        file_path: str,
        image_path: str,
        width_inches: float = 5.5,
        caption: str = None,
        paragraph_index: int = None,
        position: str = "after",
    ) -> Dict[str, Any]:
        return self._doc_ops.add_image(
            file_path,
            image_path,
            width_inches,
            caption,
            paragraph_index,
            position,
        )

    def set_paragraph_style(
        self, file_path: str, paragraph_index: int, style_name: str
    ) -> Dict[str, Any]:
        return self._doc_ops.set_paragraph_style(file_path, paragraph_index, style_name)

    def set_page_layout(
        self,
        file_path: str,
        orientation: str = None,
        margins: Dict[str, float] = None,
        header_text: str = None,
        page_numbers: bool = False,
        page_size: str = None,
    ) -> Dict[str, Any]:
        return self._doc_ops.set_page_layout(
            file_path,
            orientation,
            margins,
            header_text,
            page_numbers,
            page_size,
        )

    def _parse_range(self, range_str: str):
        from openpyxl.utils.cell import range_boundaries

        if not range_str or ":" not in range_str:
            raise ValueError(
                f"Invalid range '{range_str}'. Expected format like A1:C10"
            )
        return range_boundaries(range_str)

    def _range_has_data(self, ws, range_str: str) -> bool:
        min_col, min_row, max_col, max_row = self._parse_range(range_str)
        for row in ws.iter_rows(
            min_row=min_row,
            max_row=max_row,
            min_col=min_col,
            max_col=max_col,
            values_only=True,
        ):
            for cell in row:
                if cell is not None and str(cell).strip() != "":
                    return True
        return False

    def _resolve_formula_value(
        self, ws, expr: str, row_hint: int = None, _depth: int = 0
    ):
        if _depth > 5:
            return None

        def cell_numeric(cell_ref: str):
            val = ws[cell_ref].value
            if isinstance(val, (int, float)):
                return float(val)
            if isinstance(val, str):
                if val.startswith("="):
                    return self._resolve_formula_value(
                        ws, val[1:], row_hint=row_hint, _depth=_depth + 1
                    )
                try:
                    return float(val)
                except Exception:
                    return None
            return None

        fn_pat = re.compile(
            r"\b(SUM|AVERAGE|AVG|MIN|MAX)\(([A-Z]+\d+:[A-Z]+\d+)\)",
            re.IGNORECASE,
        )

        def fn_repl(match):
            fn = match.group(1).upper()
            rng = match.group(2)
            min_col, min_row, max_col, max_row = self._parse_range(rng)
            vals = []
            for row in ws.iter_rows(
                min_row=min_row,
                max_row=max_row,
                min_col=min_col,
                max_col=max_col,
                values_only=False,
            ):
                for c in row:
                    v = c.value
                    if isinstance(v, (int, float)):
                        vals.append(float(v))
                    elif isinstance(v, str):
                        if v.startswith("="):
                            fv = self._resolve_formula_value(
                                ws, v[1:], row_hint=c.row, _depth=_depth + 1
                            )
                            if fv is not None:
                                vals.append(float(fv))
                        else:
                            try:
                                vals.append(float(v))
                            except Exception:
                                pass
            if not vals:
                return "0"
            if fn == "SUM":
                return str(sum(vals))
            if fn in ("AVERAGE", "AVG"):
                return str(sum(vals) / len(vals))
            if fn == "MIN":
                return str(min(vals))
            if fn == "MAX":
                return str(max(vals))
            return "0"

        expr = fn_pat.sub(fn_repl, expr)
        cell_pat = re.compile(r"\b([A-Z]+\d+)\b")
        expr = cell_pat.sub(lambda m: str(cell_numeric(m.group(1)) or 0), expr)

        # Safety guard: arithmetic-only after substitutions
        if not re.fullmatch(r"[0-9\s\+\-\*/\(\)\.]+", expr):
            return None

        try:
            return float(eval(expr, {"__builtins__": {}}, {}))
        except Exception:
            return None

    def _extract_formula_functions(self, formula: str):
        if not formula:
            return set()
        expr = formula[1:] if formula.startswith("=") else formula
        return set(re.findall(r"\b([A-Z][A-Z0-9_]*)\s*\(", expr.upper()))

    def _formula_depth(self, formula: str) -> int:
        expr = (
            formula[1:]
            if isinstance(formula, str) and formula.startswith("=")
            else str(formula or "")
        )
        depth = 0
        max_depth = 0
        for ch in expr:
            if ch == "(":
                depth += 1
                max_depth = max(max_depth, depth)
            elif ch == ")":
                depth = max(0, depth - 1)
        return max_depth

    def audit_spreadsheet_formulas(
        self,
        file_path: str,
        sheet_name: str = None,
        max_examples: int = 30,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.audit_spreadsheet_formulas(file_path, sheet_name, max_examples)

    def get_formatting_info(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.get_formatting_info(file_path)

    def get_document_info(self, file_path: str) -> Dict[str, Any]:
        """Get file information for any supported format"""
        try:
            from pathlib import Path
            import os

            path = Path(file_path)
            if not path.exists():
                return {"success": False, "error": f"File not found: {file_path}"}

            stat = os.stat(file_path)
            ext = path.suffix.lower()

            info = {
                "success": True,
                "file": file_path,
                "name": path.name,
                "size": stat.st_size,
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "extension": ext,
            }

            if ext == ".docx" and DOCX_AVAILABLE:
                doc = Document(file_path)
                info["type"] = "document"
                info["paragraph_count"] = len(doc.paragraphs)
                info["word_count"] = sum(
                    len(para.text.split()) for para in doc.paragraphs
                )
            elif ext == ".xlsx" and OPENPYXL_AVAILABLE:
                wb = load_workbook(file_path, read_only=True)
                info["type"] = "spreadsheet"
                info["sheets"] = wb.sheetnames
                info["sheet_count"] = len(wb.sheetnames)
                wb.close()
            elif ext == ".pptx" and PPTX_AVAILABLE:
                prs = Presentation(file_path)
                info["type"] = "presentation"
                info["slide_count"] = len(prs.slides)

            return info
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ==================== SPREADSHEET OPERATIONS ====================

    def create_spreadsheet(
        self, output_path: str, sheet_name: str = "Sheet1"
    ) -> Dict[str, Any]:
        return self._xlsx_ops.create_spreadsheet(output_path, sheet_name)

    def write_spreadsheet(
        self,
        output_path: str,
        headers: List[str],
        data: List[List[Any]],
        sheet_name: str = "Sheet1",
        overwrite_workbook: bool = False,
        coerce_rows: bool = False,
        text_columns: List[str] = None,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.write_spreadsheet(output_path, headers, data, sheet_name, overwrite_workbook, coerce_rows, text_columns)

    def read_spreadsheet(
        self, file_path: str, sheet_name: str = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.read_spreadsheet(file_path, sheet_name)

    def calculate_column(
        self,
        file_path: str,
        column_letter: str,
        operation: str,
        sheet_name: str = "Sheet1",
        include_formulas: bool = False,
        strict_formula_safety: bool = False,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.calculate_column(file_path, column_letter, operation, sheet_name, include_formulas, strict_formula_safety)

    def _cell_to_float(self, value):
        if value is None:
            return None
        if isinstance(value, bool):
            return None
        if isinstance(value, (int, float)):
            return float(value)
        if isinstance(value, str):
            text = value.strip()
            if not text or text.startswith("="):
                return None
            try:
                return float(text)
            except Exception:
                return None
        return None

    def _value_matches_group(self, value, target) -> bool:
        if value is None:
            return False
        t = str(target).strip()
        v = str(value).strip()

        try:
            return float(v) == float(t)
        except Exception:
            return v.lower() == t.lower()

    def _category_value(self, value):
        if value is None:
            return None
        if isinstance(value, str):
            txt = value.strip()
            if not txt:
                return None
            try:
                n = float(txt)
                return int(n) if n.is_integer() else n
            except Exception:
                return txt
        if isinstance(value, float) and value.is_integer():
            return int(value)
        return value

    def frequencies(
        self,
        file_path: str,
        column_letter: str,
        sheet_name: str = "Sheet1",
        allowed_values: List[str] = None,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.frequencies(file_path, column_letter, sheet_name, allowed_values)

    def correlation_test(
        self,
        file_path: str,
        x_column: str,
        y_column: str,
        sheet_name: str = "Sheet1",
        method: str = "pearson",
    ) -> Dict[str, Any]:
        return self._xlsx_ops.correlation_test(file_path, x_column, y_column, sheet_name, method)

    def ttest_independent(
        self,
        file_path: str,
        value_column: str,
        group_column: str,
        group_a: str,
        group_b: str,
        sheet_name: str = "Sheet1",
        equal_var: bool = False,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.ttest_independent(file_path, value_column, group_column, group_a, group_b, sheet_name, equal_var)

    def mann_whitney_test(
        self,
        file_path: str,
        value_column: str,
        group_column: str,
        group_a: str,
        group_b: str,
        sheet_name: str = "Sheet1",
    ) -> Dict[str, Any]:
        return self._xlsx_ops.mann_whitney_test(file_path, value_column, group_column, group_a, group_b, sheet_name)

    def chi_square_test(
        self,
        file_path: str,
        row_column: str,
        col_column: str,
        sheet_name: str = "Sheet1",
        row_allowed_values: List[str] = None,
        col_allowed_values: List[str] = None,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.chi_square_test(file_path, row_column, col_column, sheet_name, row_allowed_values, col_allowed_values)

    def open_text_extract(
        self,
        file_path: str,
        column_letter: str,
        sheet_name: str = "Sheet1",
        limit: int = 20,
        min_length: int = 20,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.open_text_extract(file_path, column_letter, sheet_name, limit, min_length)

    def open_text_keywords(
        self,
        file_path: str,
        column_letter: str,
        sheet_name: str = "Sheet1",
        top_n: int = 15,
        min_word_length: int = 4,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.open_text_keywords(file_path, column_letter, sheet_name, top_n, min_word_length)

    def research_analysis_pack(
        self,
        file_path: str,
        sheet_name: str = "Sheet0",
        profile: str = "hlth3112",
        require_formula_safe: bool = False,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.research_analysis_pack(file_path, sheet_name, profile, require_formula_safe)

    def append_to_spreadsheet(
        self, file_path: str, row_data: List[Any], sheet_name: str = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.append_to_spreadsheet(file_path, row_data, sheet_name)

    def search_spreadsheet(
        self, file_path: str, search_text: str, sheet_name: str = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.search_spreadsheet(file_path, search_text, sheet_name)

    def add_formula(
        self, file_path: str, cell: str, formula: str, sheet_name: str = "Sheet1"
    ) -> Dict[str, Any]:
        return self._xlsx_ops.add_formula(file_path, cell, formula, sheet_name)

    # ==================== CHART OPERATIONS ====================

    def create_chart(
        self,
        file_path: str,
        chart_type: str,
        data_range: str,
        categories_range: str = None,
        title: str = "Chart",
        sheet_name: str = None,
        output_sheet: str = None,
        x_label: str = None,
        y_label: str = None,
        show_data_labels: bool = False,
        show_legend: bool = True,
        legend_pos: str = "right",
        colors: List[str] = None,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.create_chart(file_path, chart_type, data_range, categories_range, title, sheet_name, output_sheet, x_label, y_label, show_data_labels, show_legend, legend_pos, colors)

    def _create_chart_object(self, chart_type: str):
        """Create the appropriate chart object based on type"""
        chart_type = chart_type.lower()

        if chart_type == "bar" or chart_type == "column":
            chart = BarChart()
            chart.type = "col"  # Column chart (vertical bars)
        elif chart_type == "bar_horizontal":
            chart = BarChart()
            chart.type = "bar"  # Bar chart (horizontal bars)
        elif chart_type == "line":
            chart = LineChart()
        elif chart_type == "pie":
            chart = PieChart()
        elif chart_type == "scatter":
            chart = ScatterChart()
        else:
            raise ValueError(
                f"Unknown chart type: {chart_type}. Use: bar, column, bar_horizontal, line, pie, scatter"
            )

        return chart

    def _apply_chart_colors(self, chart, colors: List[str]):
        """Apply custom colors to chart series"""
        for i, ser in enumerate(chart.series):
            if i < len(colors):
                color = colors[i].lstrip("#")
                ser.graphicalProperties.solidFill.rgb = color

    def create_comparison_chart(
        self,
        file_path: str,
        chart_type: str,
        sheet_name: str = None,
        title: str = "Comparison Chart",
        start_row: int = 1,
        start_col: int = 1,
        num_categories: int = None,
        num_series: int = None,
        category_col: int = 1,
        value_cols: List[int] = None,
        output_cell: str = "A10",
        show_data_labels: bool = False,
        show_legend: bool = True,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.create_comparison_chart(file_path, chart_type, sheet_name, title, start_row, start_col, num_categories, num_series, category_col, value_cols, output_cell, show_data_labels, show_legend)

    def create_grade_distribution_chart(
        self,
        file_path: str,
        sheet_name: str = None,
        grade_column: str = "B",
        title: str = "Grade Distribution",
        output_cell: str = "F2",
    ) -> Dict[str, Any]:
        return self._xlsx_ops.create_grade_distribution_chart(file_path, sheet_name, grade_column, title, output_cell)

    def create_progress_chart(
        self,
        file_path: str,
        student_column: str = "A",
        grade_column: str = "B",
        sheet_name: str = None,
        title: str = "Student Progress",
        output_cell: str = "D2",
        show_data_labels: bool = True,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.create_progress_chart(file_path, student_column, grade_column, sheet_name, title, output_cell, show_data_labels)

    # ==================== PRESENTATION OPERATIONS ====================

    def create_presentation(
        self, output_path: str, title: str = "", subtitle: str = ""
    ) -> Dict[str, Any]:
        return self._pptx_ops.create_presentation(
            output_path, title=title, subtitle=subtitle
        )

    def add_slide(
        self, file_path: str, title: str, content: str = "", layout: str = "content"
    ) -> Dict[str, Any]:
        return self._pptx_ops.add_slide(
            file_path, title=title, content=content, layout=layout
        )

    def add_bullet_slide(
        self, file_path: str, title: str, bullets: str
    ) -> Dict[str, Any]:
        return self._pptx_ops.add_bullet_slide(file_path, title=title, bullets=bullets)

    def read_presentation(self, file_path: str) -> Dict[str, Any]:
        return self._pptx_ops.read_presentation(file_path)

    def add_image_slide(
        self, file_path: str, title: str, image_path: str
    ) -> Dict[str, Any]:
        return self._pptx_ops.add_image_slide(
            file_path, title=title, image_path=image_path
        )

    def add_table_slide(
        self,
        file_path: str,
        title: str,
        headers_csv: str,
        data_csv: str,
        coerce_rows: bool = False,
    ) -> Dict[str, Any]:
        return self._pptx_ops.add_table_slide(
            file_path,
            title=title,
            headers_csv=headers_csv,
            data_csv=data_csv,
            coerce_rows=coerce_rows,
        )

    # Helper methods
    def _hex_to_rgb(self, hex_color: str) -> tuple:
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    def _get_alignment(self, alignment: str) -> int:
        alignments = {
            "left": WD_ALIGN_PARAGRAPH.LEFT,
            "center": WD_ALIGN_PARAGRAPH.CENTER,
            "right": WD_ALIGN_PARAGRAPH.RIGHT,
            "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
        }
        return alignments.get(alignment.lower(), WD_ALIGN_PARAGRAPH.LEFT)

    # ==================== EXTENDED DOCUMENT OPERATIONS ====================

    def insert_paragraph(
        self, file_path: str, text: str, index: int, style: str = None
    ) -> Dict[str, Any]:
        return self._doc_ops.insert_paragraph(file_path, text, index, style)

    def delete_paragraph(
        self, file_path: str, index: int
    ) -> Dict[str, Any]:
        return self._doc_ops.delete_paragraph(file_path, index)

    def search_document(
        self, file_path: str, search_text: str, case_sensitive: bool = False
    ) -> Dict[str, Any]:
        return self._doc_ops.search_document(file_path, search_text, case_sensitive)

    def read_tables(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.read_tables(file_path)

    def add_hyperlink(
        self, file_path: str, text: str, url: str, paragraph_index: int = -1
    ) -> Dict[str, Any]:
        return self._doc_ops.add_hyperlink(file_path, text, url, paragraph_index)

    def add_page_break(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.add_page_break(file_path)

    def list_styles(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.list_styles(file_path)

    def set_metadata(
        self, file_path: str, author: str = None, title: str = None,
        subject: str = None, keywords: str = None, comments: str = None,
        category: str = None,
    ) -> Dict[str, Any]:
        return self._doc_ops.set_metadata(
            file_path, author, title, subject, keywords, comments, category
        )

    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.get_metadata(file_path)

    def inspect_hidden_data(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.inspect_hidden_data(file_path)

    def inspect_pdf_hidden_data(self, file_path: str) -> Dict[str, Any]:
        return self._pdf_ops.inspect_hidden_data(file_path)

    def audit_document_fonts(
        self,
        file_path: str,
        expected_font_name: str = None,
        expected_font_size: float = None,
    ) -> Dict[str, Any]:
        return self._doc_ops.audit_document_fonts(
            file_path,
            expected_font_name=expected_font_name,
            expected_font_size=expected_font_size,
        )

    def audit_document_images(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.audit_document_images(file_path)

    def document_preflight(
        self,
        file_path: str,
        expected_page_size: str = None,
        expected_font_name: str = None,
        expected_font_size: float = None,
    ) -> Dict[str, Any]:
        return self._doc_ops.document_preflight(
            file_path,
            expected_page_size=expected_page_size,
            expected_font_name=expected_font_name,
            expected_font_size=expected_font_size,
        )

    def sanitize_document(
        self,
        file_path: str,
        *,
        remove_comments: bool = False,
        accept_revisions: bool = False,
        clear_metadata: bool = False,
        remove_custom_xml: bool = False,
        set_remove_personal_information: bool = False,
        author: Optional[str] = None,
        title: Optional[str] = None,
        subject: Optional[str] = None,
        keywords: Optional[str] = None,
        output_path: Optional[str] = None,
    ) -> Dict[str, Any]:
        return self._doc_ops.sanitize_document(
            file_path,
            remove_comments=remove_comments,
            accept_revisions=accept_revisions,
            clear_metadata=clear_metadata,
            remove_custom_xml=remove_custom_xml,
            set_remove_personal_information=set_remove_personal_information,
            author=author,
            title=title,
            subject=subject,
            keywords=keywords,
            output_path=output_path,
        )

    def add_list(
        self, file_path: str, items: List[str], list_type: str = "bullet"
    ) -> Dict[str, Any]:
        return self._doc_ops.add_list(file_path, items, list_type)

    def word_count(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.word_count(file_path)

    # ==================== EXTENDED SPREADSHEET OPERATIONS ====================

    def cell_read(
        self, file_path: str, cell_ref: str, sheet_name: str = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.cell_read(file_path, cell_ref, sheet_name)

    def cell_write(
        self, file_path: str, cell_ref: str, value: str,
        sheet_name: str = None, as_text: bool = False,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.cell_write(file_path, cell_ref, value, sheet_name, as_text)

    def range_read(
        self, file_path: str, range_ref: str, sheet_name: str = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.range_read(file_path, range_ref, sheet_name)

    def delete_rows(
        self, file_path: str, start_row: int, count: int = 1, sheet_name: str = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.delete_rows(file_path, start_row, count, sheet_name)

    def delete_columns(
        self, file_path: str, start_col: int, count: int = 1, sheet_name: str = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.delete_columns(file_path, start_col, count, sheet_name)

    def sort_sheet(
        self, file_path: str, column_letter: str, sheet_name: str = None,
        descending: bool = False, numeric: bool = False,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.sort_sheet(file_path, column_letter, sheet_name, descending, numeric)

    def filter_rows(
        self, file_path: str, column_letter: str, operator: str, value: str,
        sheet_name: str = None,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.filter_rows(file_path, column_letter, operator, value, sheet_name)

    def sheet_list(self, file_path: str) -> Dict[str, Any]:
        return self._xlsx_ops.sheet_list(file_path)

    def sheet_add(
        self, file_path: str, sheet_name: str, position: int = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.sheet_add(file_path, sheet_name, position)

    def sheet_delete(self, file_path: str, sheet_name: str) -> Dict[str, Any]:
        return self._xlsx_ops.sheet_delete(file_path, sheet_name)

    def sheet_rename(
        self, file_path: str, old_name: str, new_name: str
    ) -> Dict[str, Any]:
        return self._xlsx_ops.sheet_rename(file_path, old_name, new_name)

    def merge_cells(
        self, file_path: str, range_ref: str, sheet_name: str = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.merge_cells(file_path, range_ref, sheet_name)

    def unmerge_cells(
        self, file_path: str, range_ref: str, sheet_name: str = None
    ) -> Dict[str, Any]:
        return self._xlsx_ops.unmerge_cells(file_path, range_ref, sheet_name)

    def format_cells(
        self, file_path: str, range_ref: str, sheet_name: str = None,
        bold: bool = False, italic: bool = False, font_name: str = None,
        font_size: int = None, color: str = None, bg_color: str = None,
        number_format: str = None, alignment: str = None, wrap_text: bool = False,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.format_cells(file_path, range_ref, sheet_name, bold, italic, font_name, font_size, color, bg_color, number_format, alignment, wrap_text)

    def csv_import(
        self, file_path: str, csv_path: str, sheet_name: str = "Sheet1",
        delimiter: str = ",",
    ) -> Dict[str, Any]:
        return self._xlsx_ops.csv_import(file_path, csv_path, sheet_name, delimiter)

    def csv_export(
        self, file_path: str, csv_path: str, sheet_name: str = None,
        delimiter: str = ",",
    ) -> Dict[str, Any]:
        return self._xlsx_ops.csv_export(file_path, csv_path, sheet_name, delimiter)

    # ==================== SPREADSHEET DATA VALIDATION ====================

    def add_validation(
        self, file_path: str, cell_range: str,
        validation_type: str, operator: str = None,
        formula1: str = None, formula2: str = None,
        allow_blank: bool = True, sheet_name: str = None,
        error_message: str = None, error_title: str = None,
        prompt_message: str = None, prompt_title: str = None,
        error_style: str = "stop",
    ) -> Dict[str, Any]:
        return self._xlsx_ops.add_validation(file_path, cell_range, validation_type, operator, formula1, formula2, allow_blank, sheet_name, error_message, error_title, prompt_message, prompt_title, error_style)

    def add_dropdown(
        self, file_path: str, cell_range: str, options: str,
        allow_blank: bool = True, sheet_name: str = None,
        prompt: str = None, error_message: str = None,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.add_dropdown(file_path, cell_range, options, allow_blank, sheet_name, prompt, error_message)

    def list_validations(
        self, file_path: str, sheet_name: str = None,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.list_validations(file_path, sheet_name)

    def remove_validation(
        self, file_path: str, cell_range: str = None,
        sheet_name: str = None, remove_all: bool = False,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.remove_validation(file_path, cell_range, sheet_name, remove_all)

    def validate_data(
        self, file_path: str, sheet_name: str = None,
        max_rows: int = 1000,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.validate_data(file_path, sheet_name, max_rows)

    @staticmethod
    def _check_validation(value, dv) -> tuple:
        """Check a single value against a DataValidation rule. Returns (valid, reason)."""
        vtype = dv.type
        op = dv.operator
        f1 = dv.formula1
        f2 = dv.formula2

        try:
            if vtype == "list":
                allowed = f1.strip('"').split(",") if f1 else []
                allowed = [a.strip() for a in allowed]
                if str(value).strip() not in allowed:
                    return False, f"'{value}' not in allowed list: {allowed}"
                return True, ""

            elif vtype in ("whole", "decimal"):
                try:
                    v = int(value) if vtype == "whole" else float(value)
                except (ValueError, TypeError):
                    return False, f"'{value}' is not a valid {vtype} number"
                return DocumentServerClient._check_numeric_op(v, op, f1, f2)

            elif vtype == "textLength":
                length = len(str(value))
                return DocumentServerClient._check_numeric_op(length, op, f1, f2, label="length")

            elif vtype == "date":
                # Basic date check — value should be a datetime
                from datetime import datetime as dt
                if not isinstance(value, dt):
                    return False, f"'{value}' is not a date"
                return True, ""

            elif vtype == "custom":
                # Custom formulas can't be evaluated client-side
                return True, "(custom formula — not auditable client-side)"

            else:
                return True, ""
        except Exception as e:
            return False, f"Validation error: {e}"

    @staticmethod
    def _check_numeric_op(value, operator, f1, f2, label="value"):
        """Check a numeric value against an operator and formula bounds."""
        try:
            n1 = float(f1) if f1 else None
            n2 = float(f2) if f2 else None
        except (ValueError, TypeError):
            return True, ""

        checks = {
            "between": lambda: n1 is not None and n2 is not None and n1 <= value <= n2,
            "notBetween": lambda: n1 is not None and n2 is not None and not (n1 <= value <= n2),
            "equal": lambda: n1 is not None and value == n1,
            "notEqual": lambda: n1 is not None and value != n1,
            "lessThan": lambda: n1 is not None and value < n1,
            "lessThanOrEqual": lambda: n1 is not None and value <= n1,
            "greaterThan": lambda: n1 is not None and value > n1,
            "greaterThanOrEqual": lambda: n1 is not None and value >= n1,
        }
        if operator not in checks:
            return True, ""
        if checks[operator]():
            return True, ""
        if operator == "between":
            return False, f"{label} {value} not between {n1} and {n2}"
        elif operator == "notBetween":
            return False, f"{label} {value} is between {n1} and {n2} (should not be)"
        else:
            return False, f"{label} {value} fails {operator} {n1}"

    # ==================== EXTENDED PRESENTATION OPERATIONS ====================

    def delete_slide(self, file_path: str, slide_index: int) -> Dict[str, Any]:
        return self._pptx_ops.delete_slide(file_path, slide_index=slide_index)

    def speaker_notes(
        self, file_path: str, slide_index: int, notes_text: str = None
    ) -> Dict[str, Any]:
        return self._pptx_ops.speaker_notes(
            file_path, slide_index=slide_index, notes_text=notes_text
        )

    def update_slide_text(
        self, file_path: str, slide_index: int,
        title: str = None, body: str = None,
    ) -> Dict[str, Any]:
        return self._pptx_ops.update_slide_text(
            file_path, slide_index=slide_index, title=title, body=body
        )

    def slide_count(self, file_path: str) -> Dict[str, Any]:
        return self._pptx_ops.slide_count(file_path)

    # ==================== IMAGE EXTRACTION ====================

    def extract_images_from_docx(
        self, file_path: str, output_dir: str,
        fmt: str = "png", prefix: str = "image",
    ) -> Dict[str, Any]:
        return self._doc_ops.extract_images_from_docx(
            file_path, output_dir, fmt=fmt, prefix=prefix
        )

    def extract_images_from_pptx(
        self, file_path: str, output_dir: str,
        slide_index: int = None, fmt: str = "png", prefix: str = "slide",
    ) -> Dict[str, Any]:
        return self._pptx_ops.extract_images(
            file_path,
            output_dir,
            slide_index=slide_index,
            fmt=fmt,
            prefix=prefix,
        )

    def pdf_extract_images(
        self, file_path: str, output_dir: str,
        fmt: str = "png", pages: str = None,
    ) -> Dict[str, Any]:
        return self._pdf_ops.extract_images(file_path, output_dir, fmt=fmt, pages=pages)

    @staticmethod
    def _normalize_pdf_bbox(bbox) -> Dict[str, Any]:
        return PDFOperations.normalize_bbox(bbox)

    def pdf_read_blocks(
        self,
        file_path: str,
        pages: str = None,
        include_spans: bool = True,
        include_images: bool = True,
        include_empty: bool = False,
    ) -> Dict[str, Any]:
        return self._pdf_ops.read_blocks(
            file_path,
            pages=pages,
            include_spans=include_spans,
            include_images=include_images,
            include_empty=include_empty,
        )

    def pdf_search_blocks(
        self,
        file_path: str,
        query: str,
        pages: str = None,
        case_sensitive: bool = False,
        include_spans: bool = True,
    ) -> Dict[str, Any]:
        return self._pdf_ops.search_blocks(
            file_path,
            query,
            pages=pages,
            case_sensitive=case_sensitive,
            include_spans=include_spans,
        )

    def pdf_page_to_image(
        self, file_path: str, output_dir: str,
        pages: str = None, dpi: int = 150, fmt: str = "png",
    ) -> Dict[str, Any]:
        return self._pdf_ops.page_to_image(file_path, output_dir, pages=pages, dpi=dpi, fmt=fmt)

    def pdf_sanitize(
        self,
        file_path: str,
        output_path: str = None,
        *,
        clear_metadata: bool = False,
        remove_xml_metadata: bool = False,
        author: Optional[str] = None,
        title: Optional[str] = None,
        subject: Optional[str] = None,
        keywords: Optional[str] = None,
        creator: Optional[str] = None,
        producer: Optional[str] = None,
    ) -> Dict[str, Any]:
        return self._pdf_ops.sanitize(
            file_path,
            output_path=output_path,
            clear_metadata=clear_metadata,
            remove_xml_metadata=remove_xml_metadata,
            author=author,
            title=title,
            subject=subject,
            keywords=keywords,
            creator=creator,
            producer=producer,
        )

    @staticmethod
    def _parse_page_range(pages_str: str, total: int) -> List[int]:
        return PDFOperations.parse_page_range(pages_str, total)

    def doc_render_map(self, file_path: str) -> Dict[str, Any]:
        return self._doc_ops.doc_render_map(file_path)

    # ==================== PPTX SPATIAL / TEXTBOX ====================

    def list_shapes(self, file_path: str, slide_index: int = None) -> Dict[str, Any]:
        return self._pptx_ops.list_shapes(file_path, slide_index=slide_index)

    def add_textbox(
        self, file_path: str, slide_index: int,
        text: str, left: float = 1.0, top: float = 1.0,
        width: float = 5.0, height: float = 1.5,
        font_size: float = None, font_name: str = None,
        bold: bool = False, italic: bool = False,
        color: str = None, align: str = None,
        word_wrap: bool = True,
    ) -> Dict[str, Any]:
        return self._pptx_ops.add_textbox(
            file_path,
            slide_index=slide_index,
            text=text,
            left=left,
            top=top,
            width=width,
            height=height,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            color=color,
            align=align,
            word_wrap=word_wrap,
        )

    def modify_shape(
        self, file_path: str, slide_index: int, shape_name: str,
        left: float = None, top: float = None,
        width: float = None, height: float = None,
        text: str = None, font_size: float = None,
        rotation: float = None,
    ) -> Dict[str, Any]:
        return self._pptx_ops.modify_shape(
            file_path,
            slide_index=slide_index,
            shape_name=shape_name,
            left=left,
            top=top,
            width=width,
            height=height,
            text=text,
            font_size=font_size,
            rotation=rotation,
        )

    def preview_slide(
        self, file_path: str, output_dir: str,
        slide_index: int = None, dpi: int = 150,
    ) -> Dict[str, Any]:
        return self._pptx_ops.preview_slide(
            file_path, output_dir, slide_index=slide_index, dpi=dpi
        )

    def _office_to_pdf(
        self, file_path: str, output_path: str = None,
    ) -> Dict[str, Any]:
        """Convert an OnlyOffice-supported file to PDF via x2t in Docker."""
        container_input = None
        container_pdf = None
        container_xml = None
        try:
            import subprocess
            import uuid

            abs_path = str(Path(file_path).resolve())
            if not os.path.exists(abs_path):
                return {"success": False, "error": f"File not found: {file_path}"}

            if output_path is None:
                output_path = str(Path(abs_path).with_suffix(".pdf"))
            else:
                output_path = str(Path(output_path).resolve())

            os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

            input_suffix = Path(abs_path).suffix or ".bin"
            tmp_id = uuid.uuid4().hex[:8]
            container_input = f"/tmp/convert_{tmp_id}{input_suffix}"
            container_pdf = f"/tmp/convert_{tmp_id}.pdf"
            container_xml = f"/tmp/convert_{tmp_id}.xml"

            subprocess.run(
                ["docker", "cp", abs_path, f"onlyoffice-documentserver:{container_input}"],
                check=True, capture_output=True, timeout=30,
            )

            x2t_xml = f"""<?xml version="1.0" encoding="utf-8"?>
<TaskQueueDataConvert>
  <m_sFileFrom>{container_input}</m_sFileFrom>
  <m_sFileTo>{container_pdf}</m_sFileTo>
  <m_nFormatTo>513</m_nFormatTo>
  <m_sFontDir>/usr/share/fonts</m_sFontDir>
  <m_sThemeDir>/var/www/onlyoffice/documentserver/server/FileConverter/bin/empty/themes</m_sThemeDir>
</TaskQueueDataConvert>"""
            subprocess.run(
                ["docker", "exec", "-i", "onlyoffice-documentserver", "bash", "-c",
                 f"cat > {container_xml} << 'XMLEOF'\n{x2t_xml}\nXMLEOF"],
                check=True, capture_output=True, timeout=15,
            )

            result = subprocess.run(
                ["docker", "exec", "onlyoffice-documentserver",
                 "/var/www/onlyoffice/documentserver/server/FileConverter/bin/x2t",
                 container_xml],
                capture_output=True, timeout=120,
            )
            if result.returncode != 0:
                stderr = result.stderr.decode("utf-8", errors="replace").strip()
                return {
                    "success": False,
                    "error": f"x2t conversion failed (exit {result.returncode}): {stderr}",
                }

            subprocess.run(
                ["docker", "cp", f"onlyoffice-documentserver:{container_pdf}", output_path],
                check=True, capture_output=True, timeout=30,
            )

            if not os.path.exists(output_path):
                return {"success": False, "error": "Conversion produced no output file"}

            file_size = os.path.getsize(output_path)
            pages = None
            try:
                import fitz
                doc = fitz.open(output_path)
                pages = len(doc)
                doc.close()
            except Exception:
                pass

            return {
                "success": True,
                "input_file": abs_path,
                "output_file": output_path,
                "file_size": file_size,
                "pages": pages,
            }
        except subprocess.TimeoutExpired:
            return {"success": False, "error": "Conversion timed out (is onlyoffice-documentserver running?)"}
        except FileNotFoundError:
            return {"success": False, "error": "Docker not found. Is it installed?"}
        except Exception as e:
            return {"success": False, "error": str(e)}
        finally:
            if container_input and container_pdf and container_xml:
                try:
                    subprocess.run(
                        ["docker", "exec", "onlyoffice-documentserver", "rm", "-f",
                         container_input, container_pdf, container_xml],
                        capture_output=True, timeout=10,
                    )
                except Exception:
                    pass

    def doc_to_pdf(
        self, file_path: str, output_path: str = None,
    ) -> Dict[str, Any]:
        return self._doc_ops.doc_to_pdf(file_path, output_path=output_path)

    def spreadsheet_to_pdf(
        self, file_path: str, output_path: str = None,
    ) -> Dict[str, Any]:
        return self._xlsx_ops.spreadsheet_to_pdf(file_path, output_path)

    def _preview_via_pdf(
        self,
        file_path: str,
        output_dir: str,
        converter,
        pages: str = None,
        dpi: int = 150,
        fmt: str = "png",
    ) -> Dict[str, Any]:
        """Render an Office file as page images by converting it to PDF first."""
        pdf_path = None
        try:
            os.makedirs(output_dir, exist_ok=True)
            fd, pdf_path = tempfile.mkstemp(
                prefix=f".{Path(file_path).stem}_preview_",
                suffix=".pdf",
                dir=output_dir,
            )
            os.close(fd)

            pdf_result = converter(file_path, output_path=pdf_path)
            if not pdf_result.get("success"):
                return pdf_result

            render_result = self.pdf_page_to_image(
                pdf_path, output_dir, pages=pages, dpi=dpi, fmt=fmt
            )
            if not render_result.get("success"):
                return render_result

            return {
                "success": True,
                "file": str(Path(file_path).resolve()),
                "output_dir": output_dir,
                "total_pages": render_result["total_pages"],
                "pages_rendered": render_result["pages_rendered"],
                "dpi": dpi,
                "format": fmt,
                "images": render_result["images"],
            }
        finally:
            if pdf_path:
                try:
                    os.unlink(pdf_path)
                except OSError:
                    pass

    def preview_document(
        self,
        file_path: str,
        output_dir: str,
        pages: str = None,
        dpi: int = 150,
        fmt: str = "png",
    ) -> Dict[str, Any]:
        return self._doc_ops.preview_document(
            file_path, output_dir, pages=pages, dpi=dpi, fmt=fmt
        )

    def preview_spreadsheet(
        self,
        file_path: str,
        output_dir: str,
        pages: str = None,
        dpi: int = 150,
        fmt: str = "png",
    ) -> Dict[str, Any]:
        return self._xlsx_ops.preview_spreadsheet(file_path, output_dir, pages, dpi, fmt)

    _EDITOR_EXTENSIONS = {
        "document": {".docx", ".doc", ".docm", ".dotx", ".odt", ".rtf", ".txt", ".md", ".html", ".htm", ".epub"},
        "spreadsheet": {".xlsx", ".xls", ".xlsm", ".xltx", ".xltm", ".ods", ".csv", ".tsv"},
        "presentation": {".pptx", ".ppt", ".pptm", ".ppsx", ".odp"},
        "pdf": {".pdf"},
    }

    def _editor_file_type(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        for kind, extensions in self._EDITOR_EXTENSIONS.items():
            if ext in extensions:
                return kind
        return "file"

    def _desktop_capture_tools(self) -> Dict[str, Any]:
        tools = {
            "xdotool": shutil.which("xdotool"),
            "xprop": shutil.which("xprop"),
            "onlyoffice-desktopeditors": shutil.which("onlyoffice-desktopeditors"),
            "spectacle": shutil.which("spectacle"),
            "scrot": shutil.which("scrot"),
        }
        if tools["spectacle"]:
            tools["capture_tool"] = "spectacle"
        elif tools["scrot"] and os.environ.get("XDG_SESSION_TYPE") != "wayland":
            tools["capture_tool"] = "scrot"
        else:
            tools["capture_tool"] = None
        tools["available"] = all(
            tools[name] for name in ("xdotool", "xprop", "onlyoffice-desktopeditors")
        ) and bool(tools["capture_tool"])
        return tools

    def _desktop_window_geometry(self, window_id: str) -> Dict[str, int]:
        import subprocess

        result = subprocess.run(
            ["xdotool", "getwindowgeometry", "--shell", str(window_id)],
            check=True,
            capture_output=True,
            text=True,
            timeout=10,
        )
        geometry = {}
        for line in result.stdout.splitlines():
            if "=" not in line:
                continue
            key, value = line.split("=", 1)
            if key in {"X", "Y", "WIDTH", "HEIGHT", "SCREEN"}:
                geometry[key.lower()] = int(value.strip())
        return geometry

    def _desktop_window_title(self, window_id: str) -> Dict[str, Any]:
        import subprocess

        result = subprocess.run(
            ["xprop", "-id", str(window_id), "WM_NAME", "WM_CLASS", "_NET_WM_PID"],
            check=True,
            capture_output=True,
            text=True,
            timeout=10,
        )
        title_match = re.search(r'WM_NAME\([^)]*\) = "(.*)"', result.stdout)
        class_match = re.search(r'WM_CLASS\([^)]*\) = (.*)', result.stdout)
        pid_match = re.search(r"_NET_WM_PID\(CARDINAL\) = (\d+)", result.stdout)
        classes = []
        if class_match:
            classes = re.findall(r'"([^"]+)"', class_match.group(1))
        return {
            "title": title_match.group(1) if title_match else "",
            "classes": classes,
            "pid": int(pid_match.group(1)) if pid_match else None,
        }

    def _desktop_find_editor_window(self, file_path: str) -> Dict[str, Any]:
        import subprocess

        basename = Path(file_path).name
        candidate_ids = []
        commands = [
            ["xdotool", "search", "--name", basename],
            ["xdotool", "search", "--class", "ONLYOFFICE"],
        ]
        for cmd in commands:
            result = subprocess.run(
                cmd, capture_output=True, text=True, timeout=10
            )
            if result.returncode == 0:
                candidate_ids.extend(
                    [line.strip() for line in result.stdout.splitlines() if line.strip()]
                )

        seen = set()
        candidates = []
        basename_lower = basename.lower()
        for window_id in candidate_ids:
            if window_id in seen:
                continue
            seen.add(window_id)
            try:
                meta = self._desktop_window_title(window_id)
                if "ONLYOFFICE" not in meta.get("classes", []):
                    continue
                geometry = self._desktop_window_geometry(window_id)
                title = meta.get("title", "")
                area = geometry.get("width", 0) * geometry.get("height", 0)
                candidates.append(
                    {
                        "window_id": int(window_id),
                        "title": title,
                        "pid": meta.get("pid"),
                        "classes": meta.get("classes", []),
                        "geometry": geometry,
                        "area": area,
                        "title_matches": basename_lower in title.lower(),
                    }
                )
            except Exception:
                continue

        if not candidates:
            return {}

        matched = [item for item in candidates if item["title_matches"]]
        if matched:
            candidates = matched
        else:
            return {}

        candidates.sort(
            key=lambda item: (
                0 if item["title"] else 1,
                -item["area"],
            )
        )
        return candidates[0]

    def editor_session(
        self,
        file_path: str,
        open_if_needed: bool = False,
        wait_seconds: float = 10.0,
        activate: bool = False,
    ) -> Dict[str, Any]:
        """Locate or open a native OnlyOffice Desktop Editors session for a file."""
        import subprocess
        import time

        abs_path = str(Path(file_path).resolve())
        if not os.path.exists(abs_path):
            return {"success": False, "error": f"File not found: {file_path}"}

        tools = self._desktop_capture_tools()
        if not tools["available"]:
            missing = [
                name
                for name, path in tools.items()
                if name not in {"available", "capture_tool"} and not path
            ]
            if not tools.get("capture_tool"):
                missing.append("desktop screenshot tool (spectacle or X11 scrot)")
            return {
                "success": False,
                "error": f"Desktop capture tools unavailable: {', '.join(missing)}",
            }

        session = self._desktop_find_editor_window(abs_path)
        launched = False
        if not session and open_if_needed:
            subprocess.Popen(
                [tools["onlyoffice-desktopeditors"], abs_path], start_new_session=True
            )
            launched = True
            deadline = time.monotonic() + max(1.0, float(wait_seconds))
            while time.monotonic() < deadline:
                time.sleep(0.25)
                session = self._desktop_find_editor_window(abs_path)
                if session:
                    break

        if not session:
            return {
                "success": False,
                "error": (
                    f"No OnlyOffice desktop window found for {Path(abs_path).name}. "
                    "Desktop capture needs a live editor window; use open <file> gui, "
                    "pass --open, or use --backend rendered for export-based capture."
                ),
            }

        if activate:
            self._desktop_activate_window(session["window_id"])

        geometry = session["geometry"]
        return {
            "success": True,
            "backend": "desktop",
            "file": abs_path,
            "type": self._editor_file_type(abs_path),
            "launched": launched,
            "capture_tool": tools["capture_tool"],
            "window_id": session["window_id"],
            "title": session["title"],
            "pid": session["pid"],
            "geometry": {
                "x": geometry.get("x", 0),
                "y": geometry.get("y", 0),
                "width": geometry.get("width", 0),
                "height": geometry.get("height", 0),
            },
            "supported_targets": {
                "document": ["page", "zoom"],
                "spreadsheet": ["range", "zoom"],
                "presentation": ["slide", "zoom"],
                "pdf": ["page", "zoom"],
            }.get(self._editor_file_type(abs_path), []),
        }

    def _desktop_activate_window(self, window_id: int):
        import subprocess

        subprocess.run(
            ["xdotool", "windowactivate", "--sync", str(window_id)],
            check=True,
            capture_output=True,
            timeout=10,
        )

    def _desktop_get_active_window(self) -> int:
        import subprocess

        result = subprocess.run(
            ["xdotool", "getactivewindow"],
            check=True,
            capture_output=True,
            text=True,
            timeout=10,
        )
        return int(result.stdout.strip())

    def _desktop_ensure_active_window(
        self, window_id: int, attempts: int = 3, delay_seconds: float = 0.15
    ) -> bool:
        import time

        expected = int(window_id)
        for _ in range(max(1, int(attempts))):
            self._desktop_activate_window(expected)
            time.sleep(max(0.0, float(delay_seconds)))
            try:
                if self._desktop_get_active_window() == expected:
                    return True
            except Exception:
                continue
        return False

    def _desktop_send_key(self, window_id: int, keyspec: str):
        import subprocess

        subprocess.run(
            ["xdotool", "key", "--clearmodifiers", "--window", str(window_id), keyspec],
            check=True,
            capture_output=True,
            timeout=10,
        )

    def _desktop_type_text(self, window_id: int, text: str):
        import subprocess

        subprocess.run(
            ["xdotool", "type", "--delay", "10", "--window", str(window_id), text],
            check=True,
            capture_output=True,
            timeout=10,
        )

    def _desktop_apply_viewport(
        self,
        session: Dict[str, Any],
        file_type: str,
        page: int = None,
        cell_range: str = None,
        slide: int = None,
        zoom_reset: bool = False,
        zoom_in_steps: int = 0,
        zoom_out_steps: int = 0,
        settle_ms: int = 800,
    ) -> List[str]:
        import time

        window_id = session["window_id"]
        self._desktop_activate_window(window_id)
        time.sleep(0.2)
        actions = []

        if zoom_reset:
            self._desktop_send_key(window_id, "ctrl+0")
            actions.append("zoom_reset")

        if file_type in {"document", "pdf"} and page is not None:
            self._desktop_send_key(window_id, "ctrl+g")
            time.sleep(0.15)
            self._desktop_type_text(window_id, str(int(page) + 1))
            self._desktop_send_key(window_id, "Return")
            actions.append(f"page={page}")
        elif file_type == "spreadsheet" and cell_range:
            self._desktop_send_key(window_id, "ctrl+g")
            time.sleep(0.15)
            self._desktop_type_text(window_id, cell_range)
            self._desktop_send_key(window_id, "Return")
            actions.append(f"range={cell_range}")
        elif file_type == "presentation" and slide is not None:
            self._desktop_send_key(window_id, "Home")
            if int(slide) > 0:
                subprocess.run(
                    ["xdotool", "key", "--clearmodifiers", "--window", str(window_id), "--repeat", str(int(slide)), "Page_Down"],
                    check=True,
                    capture_output=True,
                    timeout=15,
                )
            actions.append(f"slide={slide}")

        for _ in range(max(0, int(zoom_in_steps))):
            self._desktop_send_key(window_id, "ctrl+plus")
        if zoom_in_steps:
            actions.append(f"zoom_in_steps={int(zoom_in_steps)}")

        for _ in range(max(0, int(zoom_out_steps))):
            self._desktop_send_key(window_id, "ctrl+minus")
        if zoom_out_steps:
            actions.append(f"zoom_out_steps={int(zoom_out_steps)}")

        time.sleep(max(0.0, int(settle_ms)) / 1000.0)
        return actions

    def _crop_image(self, source_path: str, output_path: str, crop: str = None, fmt: str = "png") -> Dict[str, Any]:
        from PIL import Image

        img = Image.open(source_path)
        crop_box = None
        if crop:
            parts = [int(part.strip()) for part in crop.split(",")]
            if len(parts) != 4:
                raise ValueError("crop must be x,y,width,height")
            x, y, width, height = parts
            if width <= 0 or height <= 0:
                raise ValueError("crop width and height must be positive")
            crop_box = (x, y, x + width, y + height)
            img = img.crop(crop_box)
        save_format = "JPEG" if fmt.lower() in {"jpg", "jpeg"} else "PNG"
        if save_format == "JPEG" and img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        img.save(output_path, format=save_format, quality=95)
        return {
            "output_file": output_path,
            "width": img.width,
            "height": img.height,
            "crop": crop_box,
            "format": fmt.lower(),
            "size_bytes": os.path.getsize(output_path),
        }

    def _rendered_capture(
        self,
        file_path: str,
        output_path: str,
        file_type: str,
        page: int = None,
        slide: int = None,
        crop: str = None,
        dpi: int = 150,
        fmt: str = "png",
    ) -> Dict[str, Any]:
        with tempfile.TemporaryDirectory(prefix="onlyoffice-rendered-capture-") as tmpdir:
            if file_type == "document":
                preview = self.preview_document(
                    file_path, tmpdir, pages=str(page if page is not None else 0), dpi=dpi, fmt=fmt
                )
            elif file_type == "spreadsheet":
                preview = self.preview_spreadsheet(
                    file_path, tmpdir, pages=str(page if page is not None else 0), dpi=dpi, fmt=fmt
                )
            elif file_type == "presentation":
                preview = self.preview_slide(
                    file_path, tmpdir, slide_index=slide if slide is not None else 0, dpi=dpi
                )
            elif file_type == "pdf":
                preview = self.pdf_page_to_image(
                    file_path, tmpdir, pages=str(page if page is not None else 0), dpi=dpi, fmt=fmt
                )
            else:
                return {"success": False, "error": f"Rendered capture unsupported for type: {file_type}"}

            if not preview.get("success"):
                return preview
            images = preview.get("images", [])
            if not images:
                return {"success": False, "error": "Preview produced no images"}
            source_image = images[0]["file"]
            cropped = self._crop_image(source_image, output_path, crop=crop, fmt=fmt)
            return {
                "success": True,
                "backend": "rendered",
                "file": str(Path(file_path).resolve()),
                "type": file_type,
                "output_file": cropped["output_file"],
                "width": cropped["width"],
                "height": cropped["height"],
                "format": cropped["format"],
                "size_bytes": cropped["size_bytes"],
                "crop": cropped["crop"],
                "exact_viewport": False,
                "note": "Rendered fallback uses page/slide export rather than the live editor viewport.",
            }

    def capture_editor_view(
        self,
        file_path: str,
        output_path: str,
        backend: str = "auto",
        open_if_needed: bool = True,
        page: int = None,
        cell_range: str = None,
        slide: int = None,
        zoom_reset: bool = False,
        zoom_in_steps: int = 0,
        zoom_out_steps: int = 0,
        crop: str = None,
        settle_ms: int = 800,
        wait_seconds: float = 10.0,
        dpi: int = 150,
        fmt: str = None,
    ) -> Dict[str, Any]:
        """Capture either the live desktop editor viewport or a rendered fallback image."""
        abs_path = str(Path(file_path).resolve())
        if not os.path.exists(abs_path):
            return {"success": False, "error": f"File not found: {file_path}"}

        file_type = self._editor_file_type(abs_path)
        if file_type == "file":
            return {"success": False, "error": f"Unsupported file type for capture: {Path(abs_path).suffix}"}

        output_path = str(Path(output_path).resolve())
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        fmt = (fmt or Path(output_path).suffix.lstrip(".") or "png").lower()
        if fmt == "jpeg":
            fmt = "jpg"

        requested_backend = backend
        selected_backend = backend
        if selected_backend == "auto":
            selected_backend = "desktop" if self._desktop_capture_tools()["available"] else "rendered"

        if selected_backend == "rendered":
            if cell_range:
                return {
                    "success": False,
                    "error": "Rendered capture does not support live spreadsheet range targeting. Use backend=desktop for exact viewport capture.",
                }
            return self._rendered_capture(
                abs_path,
                output_path,
                file_type=file_type,
                page=page,
                slide=slide,
                crop=crop,
                dpi=dpi,
                fmt=fmt,
            )

        if selected_backend != "desktop":
            return {"success": False, "error": f"Unsupported backend: {backend}"}

        session = self.editor_session(
            abs_path, open_if_needed=open_if_needed, wait_seconds=wait_seconds, activate=True
        )
        if not session.get("success"):
            if requested_backend == "auto":
                if cell_range:
                    return {
                        "success": False,
                        "error": "Desktop session unavailable and rendered fallback cannot honor live spreadsheet range targeting. Use --open for desktop capture or omit --range.",
                    }
                return self._rendered_capture(
                    abs_path,
                    output_path,
                    file_type=file_type,
                    page=page,
                    slide=slide,
                    crop=crop,
                    dpi=dpi,
                    fmt=fmt,
                )
            return session

        actions = self._desktop_apply_viewport(
            session,
            file_type=file_type,
            page=page,
            cell_range=cell_range,
            slide=slide,
            zoom_reset=zoom_reset,
            zoom_in_steps=zoom_in_steps,
            zoom_out_steps=zoom_out_steps,
            settle_ms=settle_ms,
        )

        geometry = session["geometry"]
        if geometry["width"] <= 0 or geometry["height"] <= 0:
            return {"success": False, "error": "Editor window has invalid geometry for capture"}

        with tempfile.NamedTemporaryFile(
            prefix="onlyoffice-window-", suffix=".png", delete=False
        ) as tmp:
            tmp_path = tmp.name
        try:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
            import subprocess
            tools = self._desktop_capture_tools()
            if not self._desktop_ensure_active_window(session["window_id"]):
                return {
                    "success": False,
                    "error": (
                        "Unable to confirm the target OnlyOffice window is active for "
                        "desktop capture. Re-focus the editor and retry, or use "
                        "--backend rendered for export-based capture."
                    ),
                }
            if tools["capture_tool"] == "spectacle":
                subprocess.run(
                    ["spectacle", "-b", "-n", "-a", "-o", tmp_path],
                    check=True,
                    capture_output=True,
                    timeout=20,
                )
            elif tools["capture_tool"] == "scrot":
                subprocess.run(
                    [
                        "scrot",
                        "-a",
                        f"{max(0, geometry['x'])},{max(0, geometry['y'])},{geometry['width']},{geometry['height']}",
                        tmp_path,
                    ],
                    check=True,
                    capture_output=True,
                    timeout=20,
                )
            else:
                return {
                    "success": False,
                    "error": "No supported desktop screenshot tool is available",
                }
            cropped = self._crop_image(tmp_path, output_path, crop=crop, fmt=fmt)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

        return {
            "success": True,
            "backend": "desktop",
            "file": abs_path,
            "type": file_type,
            "output_file": cropped["output_file"],
            "width": cropped["width"],
            "height": cropped["height"],
            "format": cropped["format"],
            "size_bytes": cropped["size_bytes"],
            "crop": cropped["crop"],
            "exact_viewport": True,
            "window_id": session["window_id"],
            "title": session["title"],
            "pid": session["pid"],
            "window_geometry": geometry,
            "actions": actions,
        }

    # ==================== RDF OPERATIONS ====================
    def _rdf_format_from_path(self, file_path: str) -> str:
        return self._rdf_ops.rdf_format_from_path(file_path)

    def _rdf_safe_save(self, data: str, file_path: str) -> None:
        self._rdf_ops.rdf_safe_save(data, file_path)

    def _get_rdf_graph(self, file_path: str = None):
        return self._rdf_ops.get_graph(file_path)

    def rdf_create(
        self, file_path: str, base_uri: str = None, format: str = "turtle",
        prefixes: Dict[str, str] = None,
    ) -> Dict[str, Any]:
        return self._rdf_ops.create(
            file_path,
            base_uri=base_uri,
            format=format,
            prefixes=prefixes,
        )

    def rdf_read(
        self, file_path: str, limit: int = 100
    ) -> Dict[str, Any]:
        return self._rdf_ops.read(file_path, limit=limit)

    def rdf_add(
        self, file_path: str, subject: str, predicate: str, object_val: str,
        object_type: str = "uri", lang: str = None, datatype: str = None,
        format: str = None,
    ) -> Dict[str, Any]:
        return self._rdf_ops.add(
            file_path,
            subject,
            predicate,
            object_val,
            object_type=object_type,
            lang=lang,
            datatype=datatype,
            format=format,
        )

    def rdf_remove(
        self, file_path: str, subject: str = None, predicate: str = None,
        object_val: str = None, object_type: str = "uri", lang: str = None,
        datatype: str = None, format: str = None,
    ) -> Dict[str, Any]:
        return self._rdf_ops.remove(
            file_path,
            subject=subject,
            predicate=predicate,
            object_val=object_val,
            object_type=object_type,
            lang=lang,
            datatype=datatype,
            format=format,
        )

    def rdf_query(
        self, file_path: str, sparql: str, limit: int = 100
    ) -> Dict[str, Any]:
        return self._rdf_ops.query(file_path, sparql, limit=limit)

    def rdf_export(
        self, file_path: str, output_path: str, output_format: str = "turtle"
    ) -> Dict[str, Any]:
        return self._rdf_ops.export(
            file_path,
            output_path,
            output_format=output_format,
        )

    def rdf_merge(
        self, file_path: str, other_path: str, output_path: str = None,
        format: str = None,
    ) -> Dict[str, Any]:
        return self._rdf_ops.merge(
            file_path,
            other_path,
            output_path=output_path,
            format=format,
        )

    def rdf_stats(self, file_path: str) -> Dict[str, Any]:
        return self._rdf_ops.stats(file_path)

    def rdf_namespace(
        self, file_path: str, prefix: str = None, uri: str = None,
        format: str = None,
    ) -> Dict[str, Any]:
        return self._rdf_ops.namespace(
            file_path,
            prefix=prefix,
            uri=uri,
            format=format,
        )

    def rdf_validate(self, file_path: str, shapes_path: str) -> Dict[str, Any]:
        return self._rdf_ops.validate(file_path, shapes_path)

    # ==================== HELPERS ====================

    def check_health(self) -> bool:
        try:
            response = requests.get(f"{self.server_url}/healthcheck", timeout=5)
            return response.text.strip() == "true"
        except Exception:
            return False


_client = None


def get_client():
    global _client
    if _client is None:
        _client = DocumentServerClient(
            server_url=os.environ.get("ONLYOFFICE_URL", "http://localhost:8080"),
            secret=os.environ.get("ONLYOFFICE_SECRET", "sloane-os-secret-key-2026"),
        )
    return _client
