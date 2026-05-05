import io
import json
import os
import tempfile
import unittest
from contextlib import redirect_stdout
from unittest import mock

from PIL import Image
from pptx import Presentation

from cli_anything.onlyoffice.core import cli as cli_module
from cli_anything.onlyoffice.utils.docserver import get_client
from cli_anything.onlyoffice.utils.pptx_ops import PPTXOperations


class OnlyOfficePPTXTests(unittest.TestCase):
    def setUp(self):
        self.client = get_client()
        self.tmpdir = tempfile.TemporaryDirectory(prefix="oo-pptx-test-")
        self.base = self.tmpdir.name

    def tearDown(self):
        self.tmpdir.cleanup()

    def _path(self, name: str) -> str:
        return os.path.join(self.base, name)

    def _build_png(self, path: str, color: str = "navy") -> None:
        Image.new("RGB", (120, 60), color=color).save(path)

    def test_pptx_create_add_read_and_slide_count(self):
        path = self._path("deck.pptx")

        created = self.client.create_presentation(path, "Launch Deck", "Kick-off")
        added = self.client.add_slide(path, "Agenda", "Overview", layout="content")
        bullets = self.client.add_bullet_slide(path, "Key Points", "Alpha\nBeta\nGamma")
        read = self.client.read_presentation(path)
        count = self.client.slide_count(path)

        self.assertTrue(created["success"])
        self.assertTrue(added["success"])
        self.assertTrue(bullets["success"])
        self.assertEqual(bullets["bullets"], 3)

        self.assertTrue(read["success"])
        self.assertEqual(read["slide_count"], 3)
        self.assertEqual(read["slides"][0]["title"], "Launch Deck")
        self.assertEqual(read["slides"][1]["title"], "Agenda")
        self.assertIn("Overview", read["slides"][1]["content"][0])
        self.assertEqual(read["slides"][2]["title"], "Key Points")
        self.assertIn("Alpha", read["slides"][2]["content"][0])

        self.assertTrue(count["success"])
        self.assertEqual(count["slide_count"], 3)
        self.assertEqual(count["slide_titles"], ["Launch Deck", "Agenda", "Key Points"])

    def test_pptx_notes_update_delete_and_table_slide(self):
        path = self._path("edit_deck.pptx")
        self.assertTrue(self.client.create_presentation(path, "Deck", "Intro")["success"])
        self.assertTrue(self.client.add_slide(path, "Original Title", "Original body")["success"])

        note_write = self.client.speaker_notes(path, 1, notes_text="Remember this section")
        note_read = self.client.speaker_notes(path, 1)
        updated = self.client.update_slide_text(
            path,
            1,
            title="Updated Title",
            body="Updated body",
        )
        table = self.client.add_table_slide(
            path,
            "Results",
            "Name,Score",
            "Alice,90;Bob,88",
        )
        deleted = self.client.delete_slide(path, 2)
        count = self.client.slide_count(path)

        self.assertTrue(note_write["success"])
        self.assertTrue(note_read["success"])
        self.assertEqual(note_read["notes"], "Remember this section")
        self.assertTrue(updated["success"])
        self.assertIn("title", updated["updated_fields"])
        self.assertIn("body", updated["updated_fields"])
        self.assertTrue(table["success"])
        self.assertEqual(table["rows"], 2)
        self.assertTrue(deleted["success"])
        self.assertEqual(count["slide_count"], 2)

        presentation = Presentation(path)
        slide = presentation.slides[1]
        self.assertEqual(slide.shapes.title.text, "Updated Title")
        body_texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape != slide.shapes.title
        ]
        self.assertTrue(any("Updated body" in text for text in body_texts))

    def test_pptx_image_extraction_list_shapes_and_modify_textbox(self):
        path = self._path("shapes.pptx")
        image_path = self._path("source.png")
        out_dir = self._path("images")
        self._build_png(image_path)

        self.assertTrue(self.client.create_presentation(path, "Deck", "Intro")["success"])
        self.assertTrue(self.client.add_image_slide(path, "Evidence", image_path)["success"])

        extracted = self.client.extract_images_from_pptx(path, out_dir, fmt="png")
        before = self.client.list_shapes(path, slide_index=1)
        textbox = self.client.add_textbox(
            path,
            1,
            "Callout",
            left=1.5,
            top=1.75,
            width=3.0,
            height=1.0,
            font_size=18,
            font_name="Georgia",
            bold=True,
            italic=True,
            color="112233",
            align="center",
        )
        modified = self.client.modify_shape(
            path,
            1,
            textbox["shape_name"],
            left=2.0,
            top=2.25,
            text="Updated callout",
            font_size=20,
            rotation=15,
        )
        after = self.client.list_shapes(path, slide_index=1)

        self.assertTrue(extracted["success"])
        self.assertGreaterEqual(extracted["images_extracted"], 1)
        self.assertTrue(os.path.exists(extracted["images"][0]["file"]))

        self.assertTrue(before["success"])
        self.assertGreaterEqual(before["slides"][0]["shape_count"], 2)
        self.assertTrue(textbox["success"])
        self.assertTrue(modified["success"])

        textbox_shape = next(
            shape
            for shape in after["slides"][0]["shapes"]
            if shape["name"] == textbox["shape_name"]
        )
        self.assertEqual(textbox_shape["text"], "Updated callout")
        self.assertEqual(textbox_shape["left_inches"], 2.0)
        self.assertEqual(textbox_shape["top_inches"], 2.25)
        self.assertEqual(textbox_shape["rotation"], 15.0)

    def test_pptx_extract_images_rejects_path_prefix_and_bad_format(self):
        path = self._path("unsafe_prefix.pptx")
        image_path = self._path("source.png")
        out_dir = self._path("prefix_out")
        self._build_png(image_path)

        self.assertTrue(self.client.create_presentation(path, "Deck", "Intro")["success"])
        self.assertTrue(self.client.add_image_slide(path, "Evidence", image_path)["success"])

        traversal = self.client.extract_images_from_pptx(
            path,
            out_dir,
            fmt="png",
            prefix="../escape",
        )
        absolute = self.client.extract_images_from_pptx(
            path,
            out_dir,
            fmt="png",
            prefix=self._path("absolute"),
        )
        bad_format = self.client.extract_images_from_pptx(
            path,
            out_dir,
            fmt="../escape",
            prefix="slide",
        )

        self.assertFalse(traversal["success"])
        self.assertEqual(traversal["error_code"], "unsafe_output_prefix")
        self.assertFalse(os.path.exists(self._path("escape_01_000.png")))

        self.assertFalse(absolute["success"])
        self.assertEqual(absolute["error_code"], "unsafe_output_prefix")

        self.assertFalse(bad_format["success"])
        self.assertEqual(bad_format["error_code"], "unsupported_image_format")

    def test_pptx_extract_images_enforces_embedded_image_resource_limits(self):
        path = self._path("bounded_images.pptx")
        image_path = self._path("bounded_source.png")
        self._build_png(image_path)

        self.assertTrue(self.client.create_presentation(path, "Deck", "Intro")["success"])
        self.assertTrue(self.client.add_image_slide(path, "Evidence", image_path)["success"])

        with mock.patch.object(PPTXOperations, "MAX_EXTRACT_IMAGES", 0):
            count_limited = self.client.extract_images_from_pptx(
                path,
                self._path("count_limited"),
                fmt="png",
            )

        self.assertTrue(count_limited["success"])
        self.assertTrue(count_limited["truncated"])
        self.assertEqual(count_limited["images_extracted"], 0)
        self.assertEqual(count_limited["resource_limits"]["max_images"], 0)
        self.assertTrue(any("Stopped after 0 images" in warning for warning in count_limited["warnings"]))

        with mock.patch.object(PPTXOperations, "MAX_EXTRACT_IMAGE_COMPRESSED_BYTES", 1):
            byte_limited = self.client.extract_images_from_pptx(
                path,
                self._path("byte_limited"),
                fmt="png",
            )

        self.assertTrue(byte_limited["success"])
        self.assertEqual(byte_limited["images_extracted"], 0)
        self.assertEqual(byte_limited["images_skipped"], 1)
        self.assertTrue(byte_limited["images"][0]["skipped"])
        self.assertEqual(byte_limited["images"][0]["error_code"], "image_compressed_bytes_limit_exceeded")

        with mock.patch.object(PPTXOperations, "MAX_EXTRACT_IMAGE_PIXELS", 1):
            with mock.patch.object(Image.Image, "save") as image_save:
                pixel_limited = self.client.extract_images_from_pptx(
                    path,
                    self._path("pixel_limited"),
                    fmt="png",
                )

        self.assertTrue(pixel_limited["success"])
        self.assertEqual(pixel_limited["images_extracted"], 0)
        self.assertEqual(pixel_limited["images_skipped"], 1)
        self.assertEqual(pixel_limited["images"][0]["error_code"], "image_pixel_limit_exceeded")
        image_save.assert_not_called()

        with mock.patch.object(PPTXOperations, "MAX_EXTRACT_TOTAL_COMPRESSED_BYTES", 1):
            aggregate_byte_limited = self.client.extract_images_from_pptx(
                path,
                self._path("aggregate_byte_limited"),
                fmt="png",
            )

        self.assertTrue(aggregate_byte_limited["success"])
        self.assertTrue(aggregate_byte_limited["truncated"])
        self.assertEqual(aggregate_byte_limited["images_extracted"], 0)
        self.assertEqual(aggregate_byte_limited["images_skipped"], 1)
        self.assertEqual(
            aggregate_byte_limited["images"][0]["error_code"],
            "aggregate_image_compressed_bytes_limit_exceeded",
        )
        self.assertEqual(
            aggregate_byte_limited["resource_limits"]["max_total_compressed_image_bytes"],
            1,
        )

        with mock.patch.object(PPTXOperations, "MAX_EXTRACT_TOTAL_PIXELS", 1):
            with mock.patch.object(Image.Image, "save") as image_save:
                aggregate_pixel_limited = self.client.extract_images_from_pptx(
                    path,
                    self._path("aggregate_pixel_limited"),
                    fmt="png",
                )

        self.assertTrue(aggregate_pixel_limited["success"])
        self.assertTrue(aggregate_pixel_limited["truncated"])
        self.assertEqual(aggregate_pixel_limited["images_extracted"], 0)
        self.assertEqual(aggregate_pixel_limited["images_skipped"], 1)
        self.assertEqual(
            aggregate_pixel_limited["images"][0]["error_code"],
            "aggregate_image_pixel_limit_exceeded",
        )
        image_save.assert_not_called()

    def test_pptx_preview_uses_shared_pdf_render_pipeline(self):
        path = self._path("preview_deck.pptx")
        self.assertTrue(self.client.create_presentation(path, "Deck", "Intro")["success"])
        output_dir = self._path("preview")

        captured = {}

        def fake_office_to_pdf(file_path, output_path=None):
            self.assertEqual(file_path, path)
            self.assertIsNotNone(output_path)
            with open(output_path, "wb") as handle:
                handle.write(b"%PDF-1.4 fake")
            captured["pdf_path"] = output_path
            return {
                "success": True,
                "input_file": file_path,
                "output_file": output_path,
                "pages": 3,
            }

        def fake_pdf_page_to_image(file_path, render_dir, pages=None, dpi=150, fmt="png"):
            self.assertEqual(file_path, captured["pdf_path"])
            self.assertEqual(render_dir, output_dir)
            self.assertEqual(pages, "1")
            self.assertEqual(dpi, 200)
            self.assertEqual(fmt, "png")
            page_path = os.path.join(render_dir, "page_001.png")
            self._build_png(page_path, color="green")
            return {
                "success": True,
                "total_pages": 3,
                "pages_rendered": 1,
                "images": [
                    {
                        "page": 1,
                        "file": page_path,
                        "width": 120,
                        "height": 60,
                        "dpi": 200,
                    }
                ],
            }

        with mock.patch.object(self.client, "_office_to_pdf", side_effect=fake_office_to_pdf):
            with mock.patch.object(
                self.client, "pdf_page_to_image", side_effect=fake_pdf_page_to_image
            ):
                result = self.client.preview_slide(path, output_dir, slide_index=1, dpi=200)

        self.assertTrue(result["success"])
        self.assertEqual(result["total_slides"], 3)
        self.assertEqual(result["slides_rendered"], 1)
        self.assertEqual(result["images"][0]["slide"], 1)
        self.assertTrue(result["images"][0]["file"].endswith("slide_001.png"))
        self.assertTrue(os.path.exists(result["images"][0]["file"]))
        self.assertFalse(os.path.exists(captured["pdf_path"]))

    def test_cli_pptx_create_dispatches_via_pptx_handler(self):
        path = self._path("dispatch_create.pptx")
        stdout = io.StringIO()

        with mock.patch.object(
            cli_module.doc_server,
            "create_presentation",
            return_value={"success": True, "file": path, "slides": 1},
        ) as create_presentation:
            with mock.patch(
                "sys.argv",
                [
                    "cli-anything-onlyoffice",
                    "pptx-create",
                    path,
                    "Deck",
                    "Subtitle",
                    "--json",
                ],
            ):
                with redirect_stdout(stdout):
                    cli_module.main()

        payload = json.loads(stdout.getvalue())
        self.assertTrue(payload["success"])
        create_presentation.assert_called_once_with(path, "Deck", "Subtitle")

    def test_cli_pptx_add_textbox_dispatches_via_pptx_handler(self):
        path = self._path("dispatch_textbox.pptx")
        stdout = io.StringIO()

        with mock.patch.object(
            cli_module.doc_server,
            "add_textbox",
            return_value={"success": True, "shape_name": "TextBox 5"},
        ) as add_textbox:
            with mock.patch(
                "sys.argv",
                [
                    "cli-anything-onlyoffice",
                    "pptx-add-textbox",
                    path,
                    "1",
                    "Callout",
                    "--left",
                    "2.5",
                    "--top",
                    "1.75",
                    "--width",
                    "4.0",
                    "--height",
                    "1.2",
                    "--font-size",
                    "18",
                    "--font-name",
                    "Georgia",
                    "--bold",
                    "--italic",
                    "--color",
                    "112233",
                    "--align",
                    "center",
                    "--json",
                ],
            ):
                with redirect_stdout(stdout):
                    cli_module.main()

        payload = json.loads(stdout.getvalue())
        self.assertTrue(payload["success"])
        add_textbox.assert_called_once_with(
            path,
            1,
            "Callout",
            left=2.5,
            top=1.75,
            width=4.0,
            height=1.2,
            font_size=18.0,
            font_name="Georgia",
            bold=True,
            italic=True,
            color="112233",
            align="center",
        )
